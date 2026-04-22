"""
Figure: Translational Stratification Theory — 3-Stage Table Format
Rows: Referential Drift | Infrastructural Embedding | Stratified Legibility
Core claim: Labels travel, referents drift, infrastructures stabilize inequality.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, left, top, width, height, fill, border, radius=6, bw=Pt(1)):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = bw
    sp = sh._element
    pg = sp.find('.//a:prstGeom', NSMAP)
    if pg is not None:
        av = pg.find('a:avLst', NSMAP)
        if av is None:
            av = etree.SubElement(pg, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for c in list(av): av.remove(c)
        gd = etree.SubElement(av, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj'); gd.set('fmla', f'val {radius * 1000}')
    sh.text_frame.clear()
    return sh

def add_text(slide, left, top, width, height, text, size=10, bold=False, color=None,
             align=PP_ALIGN.LEFT, font='Calibri', italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.bold = bold; p.font.name = font; p.alignment = align
    p.font.italic = italic
    if color: p.font.color.rgb = color
    return tb

def rich_cell(slide, x, y, w, h, bg, bd, lines, bw=Pt(1), ml=0.12, mt=0.1, anchor='top'):
    """Card with multiple styled text lines."""
    card = add_rect(slide, x, y, w, h, bg, bd, radius=4, bw=bw)
    tf = card.text_frame
    tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Inches(ml); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(0.06)
    if anchor == 'middle':
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        # Set vertical centering via XML
        bodyPr = tf._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
    first = True
    for ln in lines:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln.get('t', '')
        r.font.size = Pt(ln.get('s', 10))
        r.font.bold = ln.get('b', False)
        r.font.italic = ln.get('i', False)
        r.font.name = 'Calibri'
        if 'c' in ln: r.font.color.rgb = hex_rgb(ln['c'])
        p.space_before = Pt(ln.get('sb', 1))
        p.space_after = Pt(ln.get('sa', 1))
        p.alignment = ln.get('a', PP_ALIGN.LEFT)
    return card

def add_arrow(slide, x1, y1, x2, y2, color_hex, w=2.0, dashed=False, head=True):
    cxn = slide.shapes._spTree.makeelement(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cxnSp', {})
    nv = etree.SubElement(cxn, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvCxnSpPr')
    cp = etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    cp.set('id', str(200+len(slide.shapes))); cp.set('name', f'Arr{len(slide.shapes)}')
    etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvCxnSpPr')
    etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
    sp = etree.SubElement(cxn, '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    xf = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
    if x2 < x1: xf.set('flipH', '1')
    if y2 < y1: xf.set('flipV', '1')
    off = etree.SubElement(xf, '{http://schemas.openxmlformats.org/drawingml/2006/main}off')
    off.set('x', str(int(min(x1,x2)))); off.set('y', str(int(min(y1,y2))))
    ext = etree.SubElement(xf, '{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
    ext.set('cx', str(int(abs(x2-x1)))); ext.set('cy', str(int(abs(y2-y1))))
    pg = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    pg.set('prst', 'line')
    etree.SubElement(pg, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
    c = color_hex.lstrip('#')
    ln = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    ln.set('w', str(int(w * 12700)))
    sf = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    sr = etree.SubElement(sf, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    sr.set('val', c.upper())
    if dashed:
        pd = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
        pd.set('val', 'dash')
    if head:
        te = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
        te.set('type', 'triangle'); te.set('w', 'lg'); te.set('len', 'lg')
    slide.shapes._spTree.append(cxn)

# ──────────────────────────────────────────────────────────────
# LAYOUT — Table grid: 3 rows × 5 columns
# ──────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Column widths: Stage | Mechanism | Jurisdiction Examples (wide) | Proposition
L_MARGIN  = Inches(0.35)
COL_GAP   = Inches(0.06)

STAGE_W   = Inches(1.5)        # Stage number + name
MECH_W    = Inches(2.5)        # Mechanism description
# 4 jurisdiction mini-columns inside the examples area
JURIS_W   = Inches(7.2)        # Total jurisdiction examples width
JURIS_SUB = (JURIS_W - 3 * COL_GAP) // 4   # Each jurisdiction cell
PROP_W    = Inches(1.5)        # Proposition label

HEADER_TOP = Inches(1.12)
HEADER_H   = Inches(0.42)
ROW_GAP    = Inches(0.06)

# Row heights — generous for readability
ROW_H      = Inches(1.55)
ROW1_TOP   = HEADER_TOP + HEADER_H + ROW_GAP

# Stage colors
STAGES = [
    {'num': '1', 'name': 'Referential\nDrift',
     'bg': '#EFF6FF', 'bd': '#93C5FD', 'tc': '#1E3A8A', 'acc': '#2563EB'},
    {'num': '2', 'name': 'Infrastructural\nEmbedding',
     'bg': '#ECFDF5', 'bd': '#6EE7B7', 'tc': '#064E3B', 'acc': '#059669'},
    {'num': '3', 'name': 'Stratified\nLegibility',
     'bg': '#FFFBEB', 'bd': '#FCD34D', 'tc': '#78350F', 'acc': '#D97706'},
]

# Mechanism descriptions
MECHANISMS = [
    'Portable governance terms\n("risk," "accountability,"\n"high-risk AI") travel globally\nbut their substantive\nreferents shift in each\nlocal context.',
    'Drifted categories are\nlocked into compliance\ninfrastructures — conformity\nassessments, standards,\nregistries. Apex nodes gain\ninterpretive authority.',
    'Compliance selects visibility.\nRegulators, auditors, providers\nbecome legible; workers,\ncommunities, end users\nbecome ghost nodes —\nstructurally peripheral.',
]

# Jurisdiction data: [EU, Brazil, India, Colorado]
JURISDICTIONS = ['EU AI Act', 'Brazil\nPL 2338', 'India\nNITI Aayog', 'Colorado\nSB 205']
JURIS_COLORS = [
    ('#1E3A8A', '#EFF6FF', '#93C5FD'),   # EU: blue
    ('#064E3B', '#ECFDF5', '#6EE7B7'),   # Brazil: green
    ('#7C2D12', '#FFF7ED', '#FDBA74'),   # India: orange
    ('#4C1D95', '#F5F3FF', '#C4B5FD'),   # Colorado: purple
]

JURIS_DATA = [
    # Row 1 — Referential Drift
    [
        '"Risk" = product\nsafety (Annex III\nenumeration)',
        '"Risk" = adaptive\nsocial harm\n(objective liability)',
        '"Risk" = ministerial\npriorities (develop-\nmental resilience)',
        '"Risk" = consequen-\ntial decisions\n(consumer harm)',
    ],
    # Row 2 — Infrastructural Embedding
    [
        'CE marking, notified\nbodies, EU database\n— strong closure',
        'ANPD oversight,\nLGPD-derived audit\nroutines',
        'Principles-based;\nno binding infra-\nstructure yet',
        'AG rulemaking,\nduty of care; no\nprivate right of action',
    ],
    # Row 3 — Stratified Legibility
    [
        'Regulators & auditors\nvisible; end users\ninvisible (CE system)',
        'ANPD provides\ncounter-translation\nchannel (unique)',
        'Informal pathways\nonly; no institutional\nchannel for affected',
        'AG sole gatekeeper;\nno civil society\nstanding',
    ],
]

# Proposition labels
PROP_LABELS = [
    {'id': 'P1', 'short': 'Labels travel,\nreferents shift'},
    {'id': 'P2', 'short': 'Categories\nharden into\nroutines'},
    {'id': 'P3', 'short': 'Infrastructure\nstabilizes\ninequality'},
]

# ──────────────────────────────────────────────────────────────
# BUILD
# ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = hex_rgb('#FFFFFF')

# ─── Title ───
add_text(slide, Inches(0.4), Inches(0.15), Inches(11), Inches(0.45),
         'Translational Stratification Theory — Three-Stage Pipeline',
         size=22, bold=True, color=hex_rgb('#111827'))

add_text(slide, Inches(0.4), Inches(0.62), Inches(12), Inches(0.32),
         'Labels travel, referents drift, infrastructures stabilize inequality',
         size=13, bold=False, color=hex_rgb('#6B7280'), italic=True)

# ─── Column Headers ───
hdr_x = L_MARGIN
hdrs = [
    (STAGE_W, 'STAGE'),
    (MECH_W, 'MECHANISM'),
    (JURIS_W, 'JURISDICTION-SPECIFIC EVIDENCE'),
    (PROP_W, 'PROPOSITION'),
]

for hw, hlabel in hdrs:
    rich_cell(slide, hdr_x, HEADER_TOP, hw, HEADER_H,
              hex_rgb('#1E293B'), hex_rgb('#0F172A'),
              [{'t': hlabel, 's': 9, 'b': True, 'c': '#F8FAFC', 'a': PP_ALIGN.CENTER}],
              bw=Pt(0), mt=0.08, anchor='middle')
    hdr_x += hw + COL_GAP

# Jurisdiction sub-headers inside the "EVIDENCE" area
juris_hdr_x = L_MARGIN + STAGE_W + COL_GAP + MECH_W + COL_GAP
juris_hdr_y = HEADER_TOP + HEADER_H + ROW_GAP
JURIS_HDR_H = Inches(0.32)

for ji, (jname, (jtxt, jbg, jbd)) in enumerate(zip(JURISDICTIONS, JURIS_COLORS)):
    jx = juris_hdr_x + ji * (JURIS_SUB + COL_GAP)
    rich_cell(slide, jx, juris_hdr_y, JURIS_SUB, JURIS_HDR_H,
              hex_rgb(jbg), hex_rgb(jbd),
              [{'t': jname, 's': 9, 'b': True, 'c': jtxt, 'a': PP_ALIGN.CENTER}],
              bw=Pt(1), mt=0.04, anchor='middle')

# Adjust row start below jurisdiction sub-headers
ROW1_TOP = juris_hdr_y + JURIS_HDR_H + ROW_GAP

# ─── Data Rows ───
for ri, stage in enumerate(STAGES):
    ry = ROW1_TOP + ri * (ROW_H + ROW_GAP)
    cx = L_MARGIN

    # ── Stage cell (number + name) ──
    rich_cell(slide, cx, ry, STAGE_W, ROW_H,
              hex_rgb(stage['bg']), hex_rgb(stage['bd']),
              [
                  {'t': stage['num'], 's': 28, 'b': True, 'c': stage['acc'], 'a': PP_ALIGN.CENTER, 'sa': 2},
                  {'t': stage['name'], 's': 12, 'b': True, 'c': stage['tc'], 'a': PP_ALIGN.CENTER, 'sb': 0},
              ],
              bw=Pt(1.5), mt=0.15, anchor='middle')
    cx += STAGE_W + COL_GAP

    # ── Mechanism cell ──
    rich_cell(slide, cx, ry, MECH_W, ROW_H,
              hex_rgb('#F8FAFC'), hex_rgb('#CBD5E1'),
              [{'t': MECHANISMS[ri], 's': 9.5, 'b': False, 'c': '#334155'}],
              mt=0.12)
    cx += MECH_W + COL_GAP

    # ── Jurisdiction cells (4 sub-columns) ──
    for ji, cell_text in enumerate(JURIS_DATA[ri]):
        jtxt, jbg, jbd = JURIS_COLORS[ji]
        jx = cx + ji * (JURIS_SUB + COL_GAP)

        # Special: ghost node highlight for row 3
        if ri == 2 and ji != 1:  # Stratified Legibility row, non-Brazil (Brazil has counter-translation)
            cell_bg = '#FFFBEB'
            cell_bd = '#FCD34D'
        elif ri == 2 and ji == 1:  # Brazil counter-translation
            cell_bg = '#FEF2F2'
            cell_bd = '#FCA5A5'
            jtxt = '#991B1B'
        else:
            cell_bg = '#FFFFFF'
            cell_bd = jbd

        cell_lines = []
        for li, line in enumerate(cell_text.split('\n')):
            cell_lines.append({
                't': line, 's': 10 if li == 0 else 9,
                'b': (li == 0), 'c': jtxt,
                'sb': 0 if li == 0 else 1, 'sa': 1,
            })

        rich_cell(slide, jx, ry, JURIS_SUB, ROW_H,
                  hex_rgb(cell_bg), hex_rgb(cell_bd), cell_lines,
                  ml=0.1, mt=0.1)

    # ── Proposition cell ──
    prop = PROP_LABELS[ri]
    prop_x = cx + JURIS_W + COL_GAP
    rich_cell(slide, prop_x, ry, PROP_W, ROW_H,
              hex_rgb('#EEF2FF'), hex_rgb('#A5B4FC'),
              [
                  {'t': prop['id'], 's': 22, 'b': True, 'c': '#4F46E5', 'a': PP_ALIGN.CENTER, 'sa': 4},
                  {'t': prop['short'], 's': 9, 'b': False, 'c': '#4338CA', 'a': PP_ALIGN.CENTER, 'i': True, 'sb': 0},
              ],
              bw=Pt(1.5), mt=0.12, anchor='middle')

# ──────────────────────────────────────────────────────────────
# FLOW ARROWS (between rows on the left)
# ──────────────────────────────────────────────────────────────

arrow_x = int(L_MARGIN + STAGE_W // 2)
for ri in range(2):
    ay1 = int(ROW1_TOP + ri * (ROW_H + ROW_GAP) + ROW_H + Inches(0.01))
    ay2 = int(ROW1_TOP + (ri + 1) * (ROW_H + ROW_GAP) - Inches(0.01))
    add_arrow(slide, arrow_x, ay1, arrow_x, ay2, '#1E293B', w=2.5)

# ──────────────────────────────────────────────────────────────
# OUTCOME BANNER — Sedimentation / Counter-Translation
# ──────────────────────────────────────────────────────────────

banner_y = ROW1_TOP + 3 * (ROW_H + ROW_GAP) + Inches(0.06)
banner_w = SLIDE_W - 2 * L_MARGIN
banner_h = Inches(0.75)

# Banner background
add_rect(slide, L_MARGIN, banner_y, banner_w, banner_h,
         hex_rgb('#F8FAFC'), hex_rgb('#CBD5E1'), radius=5, bw=Pt(1.2))

# Left half: Sedimentation
sedi_w = banner_w // 2 - Inches(0.1)
rich_cell(slide, L_MARGIN + Inches(0.08), banner_y + Inches(0.06),
          sedi_w, banner_h - Inches(0.12),
          hex_rgb('#EEF2FF'), hex_rgb('#A5B4FC'),
          [
              {'t': '⬤  Sedimentation', 's': 11, 'b': True, 'c': '#4338CA', 'sa': 2},
              {'t': 'Categories stabilize across systems → infrastructural closure → resistance shifts to recoding',
               's': 9, 'b': False, 'c': '#4338CA', 'sb': 0},
          ],
          bw=Pt(1), mt=0.06)

# Right half: Counter-Translation
ct_x = L_MARGIN + banner_w // 2 + Inches(0.1)
rich_cell(slide, ct_x, banner_y + Inches(0.06),
          sedi_w, banner_h - Inches(0.12),
          hex_rgb('#FEF2F2'), hex_rgb('#FCA5A5'),
          [
              {'t': '⬤  Counter-Translation', 's': 11, 'b': True, 'c': '#991B1B', 'sa': 2},
              {'t': 'Ghost nodes reframe terms (risk as collective harm) → reopens referential drift → feeds back to Stage 1',
               's': 9, 'b': False, 'c': '#B91C1C', 'sb': 0},
          ],
          bw=Pt(1), mt=0.06)

# Feedback arrows from outcome banner
# Sedimentation: banner → left side of row 2
sedi_cx = int(L_MARGIN + Inches(0.08) + sedi_w // 2)
add_arrow(slide, sedi_cx, int(banner_y), sedi_cx,
          int(ROW1_TOP + ROW_H + ROW_GAP + ROW_H + Inches(0.02)),
          '#4338CA', w=1.5, dashed=True)

# Counter-translation: banner → left side of row 1
ct_cx = int(ct_x + sedi_w // 2)
ct_loop_y = int(banner_y + banner_h + Inches(0.06))
ct_loop_x = int(L_MARGIN - Inches(0.12))
# Down from banner
add_arrow(slide, ct_cx, int(banner_y + banner_h), ct_cx, ct_loop_y, '#DC2626', w=1.5, dashed=True, head=False)
# Left to edge
add_arrow(slide, ct_cx, ct_loop_y, ct_loop_x, ct_loop_y, '#DC2626', w=1.5, dashed=True, head=False)
# Up to row 1
add_arrow(slide, ct_loop_x, ct_loop_y, ct_loop_x, int(ROW1_TOP + ROW_H // 2), '#DC2626', w=1.5, dashed=True)
# Right into row 1
add_arrow(slide, ct_loop_x, int(ROW1_TOP + ROW_H // 2), int(L_MARGIN), int(ROW1_TOP + ROW_H // 2),
          '#DC2626', w=1.5, dashed=True)

# ──────────────────────────────────────────────────────────────
# LEGEND
# ──────────────────────────────────────────────────────────────

leg_y = Inches(7.1)

sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), leg_y + Inches(0.06), Inches(0.4), Pt(3.5))
sh.fill.solid(); sh.fill.fore_color.rgb = hex_rgb('#1E293B'); sh.line.fill.background()
add_text(slide, Inches(1.0), leg_y - Inches(0.02), Inches(2.5), Inches(0.22),
         '→ Sequential (1 → 2 → 3)', size=8, color=hex_rgb('#374151'))

sh2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), leg_y + Inches(0.06), Inches(0.4), Pt(3.5))
sh2.fill.solid(); sh2.fill.fore_color.rgb = hex_rgb('#4338CA'); sh2.line.fill.background()
add_text(slide, Inches(4.0), leg_y - Inches(0.02), Inches(3.2), Inches(0.22),
         '⇢ Sedimentation (Stage 3 → 2)', size=8, color=hex_rgb('#374151'))

sh3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), leg_y + Inches(0.06), Inches(0.4), Pt(3.5))
sh3.fill.solid(); sh3.fill.fore_color.rgb = hex_rgb('#DC2626'); sh3.line.fill.background()
add_text(slide, Inches(7.5), leg_y - Inches(0.02), Inches(3.8), Inches(0.22),
         '⇢ Counter-translation (Stage 3 → 1)', size=8, color=hex_rgb('#374151'))

# ──────────────────────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────────────────────

output = r'c:\Users\mount\.gemini\antigravity\scratch\TST_Diagram_Publication.pptx'
prs.save(output)
print(f'[OK] Saved: {output}')
