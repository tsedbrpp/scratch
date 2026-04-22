"""
Organizational Impacts of Translational Stratification – 3‑Stage Pipeline
Creates a PowerPoint slide that visualises how the three stages (Referential Drift, Infrastructure Embedding, Stratified Legibility) reshape organisational structures: authority concentration, compliance routing, and ghost‑node marginalisation.
All shapes are native PowerPoint objects, fully editable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

# ---------------------------------------------------------------------------
# Helper utilities (reuse patterns from earlier diagram scripts)
# ---------------------------------------------------------------------------

def hex_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_rect(slide, left, top, width, height, fill, border, radius=6, bw=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border; shape.line.width = bw
    # Adjust corner radius via XML
    sp = shape._element
    prst_geom = sp.find('.//a:prstGeom', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
    if prst_geom is not None:
        av_lst = prst_geom.find('a:avLst', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        if av_lst is None:
            av_lst = etree.SubElement(prst_geom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for child in list(av_lst):
                av_lst.remove(child)
        gd = etree.SubElement(av_lst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj')
        gd.set('fmla', f'val {radius * 1000}')
    shape.text_frame.clear()
    return shape

def add_textbox(slide, left, top, width, height, text, size=10, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = 'Inter'
    p.font.italic = italic
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return tb

def rich_cell(slide, left, top, width, height, bg, bd, lines, radius=5, bw=Pt(1)):
    rect = add_rect(slide, left, top, width, height, bg, bd, radius=radius, bw=bw)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.04)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line.get('t', '')
        run.font.size = Pt(line.get('s', 10))
        run.font.bold = line.get('b', False)
        run.font.italic = line.get('i', False)
        run.font.name = 'Inter'
        if 'c' in line:
            run.font.color.rgb = hex_rgb(line['c'])
        p.alignment = line.get('a', PP_ALIGN.LEFT)
    return rect

def add_arrow(slide, x1, y1, x2, y2, color_hex, width=1.5, dashed=False, head=True):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, x1, y1, x2 - x1, y2 - y1
    )
    line = conn.line
    line.width = Pt(width)
    line.fill.solid()
    line.fill.fore_color.rgb = hex_rgb(color_hex)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if head:
        conn.end_arrowhead.style = conn.end_arrowhead.STYLE_TRIANGLE
        conn.end_arrowhead.width = conn.end_arrowhead.WIDTH_LARGE
        conn.end_arrowhead.length = conn.end_arrowhead.LENGTH_LARGE
    return conn

# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.35)
USABLE_W = SLIDE_W - 2 * MARGIN

# Columns: Stage (P1‑P3) | Impact (Authority, Routing, Ghost Nodes)
COL_STAGE = Inches(2.2)
COL_IMPACT = USABLE_W - COL_STAGE - Inches(0.1)
GAP = Inches(0.08)

HEADER_H = Inches(0.4)
ROW_H = Inches(1.1)
ROW_GAP = Inches(0.07)
START_Y = MARGIN + HEADER_H + ROW_GAP

# ---------------------------------------------------------------------------
# Data for each stage
# ---------------------------------------------------------------------------

STAGES = [
    {
        'label': 'P1 – Referential Drift',
        'bg': '#F3F4F6',
        'bd': '#D1D5DB',
        'impact': [
            {'t': 'Authority concentrates in technocratic nodes', 'c': '#111827', 's': 9},
            {'t': 'Policy vocabularies travel but acquire divergent referents', 'c': '#111827', 's': 9},
            {'t': 'Actors that matter shift (expertise, obligations)', 'c': '#111827', 's': 9},
        ]
    },
    {
        'label': 'P2 – Infrastructure Embedding',
        'bg': '#F0FDF4',
        'bd': '#86EFAC',
        'impact': [
            {'t': 'Governance routed through standards, audits, procurement', 'c': '#065F46', 's': 9},
            {'t': 'Apex institutions become obligatory passage points', 'c': '#065F46', 's': 9},
            {'t': 'Stratified legibility – compliant actors become visible', 'c': '#065F46', 's': 9},
        ]
    },
    {
        'label': 'P3 – Stratified Legibility',
        'bg': '#FEF2F2',
        'bd': '#FECACA',
        'impact': [
            {'t': 'Civil society, workers, affected communities become ghost nodes', 'c': '#991B1B', 's': 9},
            {'t': 'Sedimentation vs. counter‑translation shapes future openings', 'c': '#991B1B', 's': 9},
            {'t': 'Uneven contestation capacity reinforces existing hierarchies', 'c': '#991B1B', 's': 9},
        ]
    },
]

# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex_rgb('#FFFFFF')

# Title
add_textbox(slide, MARGIN, Inches(0.12), USABLE_W, Inches(0.35),
            'Organizational Impacts of Translational Stratification',
            size=20, bold=True, color=hex_rgb('#111827'))
add_textbox(slide, MARGIN, Inches(0.55), USABLE_W, Inches(0.3),
            'How portable governance vocabularies reshape authority, compliance pathways, and marginalised actors',
            size=10.5, color=hex_rgb('#6B7280'), italic=True)

# Header cells
rich_cell(slide, MARGIN, MARGIN, COL_STAGE, HEADER_H,
          hex_rgb('#E5E7EB'), hex_rgb('#9CA3AF'),
          [{'t': 'STAGE', 's': 10, 'b': True, 'c': '#111827', 'a': PP_ALIGN.CENTER}], radius=4)
rich_cell(slide, MARGIN + COL_STAGE + GAP, MARGIN, COL_IMPACT, HEADER_H,
          hex_rgb('#E5E7EB'), hex_rgb('#9CA3AF'),
          [{'t': 'ORGANIZATIONAL IMPACT', 's': 10, 'b': True, 'c': '#111827', 'a': PP_ALIGN.CENTER}], radius=4)

# Rows per stage
for i, st in enumerate(STAGES):
    y = START_Y + i * (ROW_H + ROW_GAP)
    # Stage cell
    rich_cell(slide, MARGIN, y, COL_STAGE, ROW_H,
              hex_rgb(st['bg']), hex_rgb(st['bd']),
              [{'t': st['label'], 's': 11, 'b': True, 'c': '#111827', 'a': PP_ALIGN.CENTER}], radius=5)
    # Impact cell – bullet list
    rich_cell(slide, MARGIN + COL_STAGE + GAP, y, COL_IMPACT, ROW_H,
              hex_rgb(st['bg']), hex_rgb(st['bd']),
              st['impact'], radius=5)
    # Dashed connector between rows (visual flow)
    if i < len(STAGES) - 1:
        add_arrow(slide,
                  int(MARGIN + COL_STAGE/2), int(y + ROW_H),
                  int(MARGIN + COL_STAGE/2), int(y + ROW_H + ROW_GAP),
                  '#6B7280', width=1.2, dashed=True, head=False)

# Bottom callout – synthesis
call_y = START_Y + len(STAGES)*(ROW_H+ROW_GAP) + Inches(0.12)
call_h = Inches(0.55)
call_bg = add_rect(slide, MARGIN, call_y, USABLE_W, call_h,
                   hex_rgb('#FEF2F2'), hex_rgb('#FECACA'), radius=5, bw=Pt(1.2))
# Accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, call_y, Inches(0.06), call_h)
bar.fill.solid(); bar.fill.fore_color.rgb = hex_rgb('#991B1B'); bar.line.fill.background()
# Callout text
tf = call_bg.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.06)
para = tf.paragraphs[0]
run = para.add_run()
run.text = ('Portable vocabularies travel, but local translation concentrates authority in technocratic nodes, routes governance through compliance infrastructures, and marginalises civil‑society, workers, and affected communities as ghost nodes. Counter‑translation can reopen pathways, yet capacity is unevenly distributed.')
run.font.size = Pt(8.5); run.font.bold = True; run.font.color.rgb = hex_rgb('#991B1B')

# Legend
legend_y = call_y + call_h + Inches(0.12)
add_textbox(slide, MARGIN, legend_y, USABLE_W, Inches(0.2),
            'Legend: Dashed vertical arrows indicate the sequential flow across stages. Colours follow the Policy Prism palette; orange callout highlights the core organizational impact.',
            size=9, color=hex_rgb('#6B7280'))

# Save PPTX
out_path = r'c:\\Users\\mount\\.gemini\\antigravity\\scratch\\Organizational_Impacts_Translational_Stratification.pptx'
prs.save(out_path)
print(f'[OK] Saved: {out_path}')
