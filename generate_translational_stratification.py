"""
Translational Stratification Analysis – Three‑Stage Pipeline Diagram
Generates a PowerPoint slide visualizing the three stages (Referential Drift, Infrastructure Embedding, Stratified Legibility)
across five jurisdictions (EU, Brazil, India, Colorado, United States).
All shapes are native PowerPoint objects – fully editable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

# ---------------------------------------------------------------------------
# Helper utilities (reuse patterns from earlier diagram scripts)
# ---------------------------------------------------------------------------

def hex_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_rect(slide, left, top, width, height, fill, border, radius=6, bw=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border; shape.line.width = bw
    # Adjust corner radius via XML
    sp = shape._element
    prst_geom = sp.find('.//a:prstGeom', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
    if prst_geom is not None:
        av_lst = prst_geom.find('a:avLst', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        if av_lst is None:
            av_lst = etree.SubElement(prst_geom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for child in list(av_lst):
                av_lst.remove(child)
        gd = etree.SubElement(av_lst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj')
        gd.set('fmla', f'val {radius * 1000}')
    shape.text_frame.clear()
    return shape

def add_textbox(slide, left, top, width, height, text, size=10, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = 'Inter'
    p.font.italic = italic
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return tb

def rich_cell(slide, left, top, width, height, bg, bd, lines, radius=5, bw=Pt(1)):
    rect = add_rect(slide, left, top, width, height, bg, bd, radius=radius, bw=bw)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.04)
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line.get('t', '')
        run.font.size = Pt(line.get('s', 10))
        run.font.bold = line.get('b', False)
        run.font.italic = line.get('i', False)
        run.font.name = 'Inter'
        if 'c' in line:
            run.font.color.rgb = hex_rgb(line['c'])
        p.alignment = line.get('a', PP_ALIGN.LEFT)
    return rect

def add_arrow(slide, x1, y1, x2, y2, color_hex, width=1.5, dashed=False, head=True):
    # Use native connector shape for simplicity
    conn = slide.shapes.add_connector(
        MSO_SHAPE.LINE, x1, y1, x2 - x1, y2 - y1
    )
    line = conn.line
    line.width = Pt(width)
    line.fill.solid()
    line.fill.fore_color.rgb = hex_rgb(color_hex)
    if dashed:
        line.dash_style = line.DASH_STYLE_DASH
    if head:
        conn.end_arrowhead.style = conn.end_arrowhead.STYLE_TRIANGLE
        conn.end_arrowhead.width = conn.end_arrowhead.WIDTH_LARGE
        conn.end_arrowhead.length = conn.end_arrowhead.LENGTH_LARGE
    return conn

# ---------------------------------------------------------------------------
# Data definition – jurisdiction specific colours and texts
# ---------------------------------------------------------------------------

JURIS = [
    {
        'name': 'European Union',
        'color': '#6366F1',
        'light': '#EEF2FF',
        'border': '#A5B4FC',
        'drift': 'Risk → product safety',
        'embed': 'CE‑mark, notified bodies, EU‑DB',
        'legible': 'Regulators/auditors visible; users invisible (Ghost Nodes)'
    },
    {
        'name': 'Brazil',
        'color': '#0D9488',
        'light': '#F0FDFA',
        'border': '#5EEAD4',
        'drift': 'Risk → adaptive social harm',
        'embed': 'ANPD oversight + LGPD audits',
        'legible': 'Formal counter‑translation (ANPD re‑classification)'
    },
    {
        'name': 'India',
        'color': '#D97706',
        'light': '#FFFBEB',
        'border': '#FCD34D',
        'drift': 'Risk → developmental resilience',
        'embed': 'Principles‑based ministerial coordination',
        'legible': 'Informal pathways only (no formal challenge)'
    },
    {
        'name': 'Colorado',
        'color': '#16A34A',
        'light': '#F0FDF4',
        'border': '#86EFAC',
        'drift': 'Risk → consequential consumer decisions',
        'embed': 'Attorney General rulemaking; no private right of action',
        'legible': 'AG as sole gatekeeper (Ghost Nodes)'
    },
    {
        'name': 'United States',
        'color': '#6B7280',
        'light': '#F9FAFB',
        'border': '#9CA3AF',
        'drift': 'Risk → security‑industrial competition',
        'embed': 'Procurement, OMB/CAISI/CAIOC, export controls, industry‑led standards',
        'legible': 'Federal/industry evaluators visible; public & workers only selectively enrolled'
    },
]

# ---------------------------------------------------------------------------
# Layout constants (slide size, margins, column widths)
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.35)
USABLE_W = SLIDE_W - 2 * MARGIN

# Column widths (proportional)
COL_JUR = Inches(1.5)
COL_DRIFT = Inches(3.2)
COL_EMBED = Inches(3.2)
COL_LEGIBLE = Inches(3.2)
GAP = Inches(0.08)

# Vertical spacing
HEADER_H = Inches(0.38)
ROW_H = Inches(1.02)
ROW_GAP = Inches(0.06)
START_Y = MARGIN + HEADER_H + ROW_GAP

# X positions
X_JUR = MARGIN
X_DRIFT = X_JUR + COL_JUR + GAP
X_EMBED = X_DRIFT + COL_DRIFT + GAP
X_LEGIBLE = X_EMBED + COL_EMBED + GAP

# ---------------------------------------------------------------------------
# Build the presentation
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])
# Background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = hex_rgb('#FFFFFF')

# Title
add_textbox(slide, MARGIN, Inches(0.12), USABLE_W, Inches(0.4),
            'Translational Stratification Analysis – Three‑Stage Pipeline',
            size=20, bold=True, color=hex_rgb('#111827'))
add_textbox(slide, MARGIN, Inches(0.52), USABLE_W, Inches(0.35),
            'Stable vocabularies travel, but local referents create inequity through layered compliance infrastructures.',
            size=10.5, color=hex_rgb('#6B7280'), italic=True)

# Header row
rich_cell(slide, X_JUR, MARGIN + HEADER_H, COL_JUR, HEADER_H,
          hex_rgb('#F3F4F6'), hex_rgb('#D1D5DB'),
          [{'t': 'JURISDICTION', 's': 9, 'b': True, 'c': '#111827', 'a': PP_ALIGN.CENTER}], radius=4, bw=Pt(0))
rich_cell(slide, X_DRIFT, MARGIN + HEADER_H, COL_DRIFT, HEADER_H,
          hex_rgb('#FFFBEB'), hex_rgb('#FDE68A'),
          [{'t': 'REFERENTIAL DRIFT', 's': 9, 'b': True, 'c': '#78350F', 'a': PP_ALIGN.CENTER}], radius=4)
rich_cell(slide, X_EMBED, MARGIN + HEADER_H, COL_EMBED, HEADER_H,
          hex_rgb('#F0FDF4'), hex_rgb('#86EFAC'),
          [{'t': 'INFRASTRUCTURE EMBEDDING', 's': 9, 'b': True, 'c': '#065F46', 'a': PP_ALIGN.CENTER}], radius=4)
rich_cell(slide, X_LEGIBLE, MARGIN + HEADER_H, COL_LEGIBLE, HEADER_H,
          hex_rgb('#F9FAFB'), hex_rgb('#D1D5DB'),
          [{'t': 'STRATIFIED LEGIBILITY', 's': 9, 'b': True, 'c': '#374151', 'a': PP_ALIGN.CENTER}], radius=4)

# Rows per jurisdiction
for idx, jur in enumerate(JURIS):
    y = START_Y + idx * (ROW_H + ROW_GAP)
    # Jurisdiction cell
    rich_cell(slide, X_JUR, y, COL_JUR, ROW_H,
              hex_rgb(jur['light']), hex_rgb(jur['border']),
              [{'t': jur['name'], 's': 11, 'b': True, 'c': jur['color'], 'a': PP_ALIGN.CENTER}], radius=5)
    # Drift cell
    rich_cell(slide, X_DRIFT, y, COL_DRIFT, ROW_H,
              hex_rgb(jur['light']), hex_rgb(jur['border']),
              [{'t': jur['drift'], 's': 9, 'c': '#78350F', 'a': PP_ALIGN.CENTER}], radius=5)
    # Embed cell
    rich_cell(slide, X_EMBED, y, COL_EMBED, ROW_H,
              hex_rgb(jur['light']), hex_rgb(jur['border']),
              [{'t': jur['embed'], 's': 9, 'c': '#065F46', 'a': PP_ALIGN.CENTER}], radius=5)
    # Legible cell
    rich_cell(slide, X_LEGIBLE, y, COL_LEGIBLE, ROW_H,
              hex_rgb(jur['light']), hex_rgb(jur['border']),
              [{'t': jur['legible'], 's': 9, 'c': '#374151', 'a': PP_ALIGN.CENTER}], radius=5)
    # Dashed arrows from drift → embed → legible (same colour)
    # Arrow from centre of drift cell to centre of embed cell
    add_arrow(slide,
              int(X_DRIFT + COL_DRIFT/2), int(y + ROW_H/2),
              int(X_EMBED + COL_EMBED/2), int(y + ROW_H/2),
              jur['color'], width=1.2, dashed=True, head=False)
    # Arrow from centre of embed cell to centre of legible cell
    add_arrow(slide,
              int(X_EMBED + COL_EMBED/2), int(y + ROW_H/2),
              int(X_LEGIBLE + COL_LEGIBLE/2), int(y + ROW_H/2),
              jur['color'], width=1.2, dashed=True, head=False)

# Bottom callout
call_y = START_Y + len(JURIS)*(ROW_H+ROW_GAP) + Inches(0.08)
call_h = Inches(0.55)
call_bg = add_rect(slide, MARGIN, call_y, USABLE_W, call_h,
                   hex_rgb('#FEF2F2'), hex_rgb('#FECACA'), radius=5, bw=Pt(1.2))
# Accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, call_y, Inches(0.06), call_h)
bar.fill.solid(); bar.fill.fore_color.rgb = hex_rgb('#991B1B'); bar.line.fill.background()
# Callout text
tf = call_bg.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.06)
para = tf.paragraphs[0]
run1 = para.add_run(); run1.text = 'Stable vocabularies travel, but local referents create inequity through layered compliance infrastructures.'
run1.font.size = Pt(8.5); run1.font.bold = True; run1.font.color.rgb = hex_rgb('#991B1B')

# Legend row (simple text)
legend_y = call_y + call_h + Inches(0.12)
add_textbox(slide, MARGIN, legend_y, USABLE_W, Inches(0.2),
            'Legend: ↦ Dashed arrows show the flow from drift → embedding → legibility. Colours correspond to jurisdictions.',
            size=9, color=hex_rgb('#6B7280'))

# Save PPTX
out_path = r'c:\Users\mount\.gemini\antigravity\scratch\Translational_Stratification_Analysis.pptx'
prs.save(out_path)
print(f'[OK] Saved: {out_path}')
