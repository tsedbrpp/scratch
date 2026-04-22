"""
Build PowerPoint with all three figures as high-fidelity rendered images.
Figure 1: Directed Accountability Flows
Figure 2: Institutional Logics & Accountability Patterns
Figure 3: Ghost Node Exclusion Matrix
"""

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import glob, os

brain_dir = r'C:\Users\mount\.gemini\antigravity\brain\62c432ec-01b0-4622-8e0c-b2a72cf05648'

def find_latest(pattern):
    candidates = glob.glob(os.path.join(brain_dir, pattern))
    if not candidates:
        raise FileNotFoundError(f'No files matching {pattern}')
    return sorted(candidates)[-1]

fig1_path = find_latest('accountability_figure_1_*.png')
fig2_path = find_latest('accountability_figure_2_*.png')
fig3_path = find_latest('ghost_node_matrix_figure_*.png')

print(f'Fig1: {fig1_path}')
print(f'Fig2: {fig2_path}')
print(f'Fig3: {fig3_path}')

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
margin = Inches(0.3)
img_width = prs.slide_width - 2 * margin
img_height = prs.slide_height - 2 * margin

for fig_path, title in [
    (fig1_path, 'Figure 1 - Directed Accountability Flows'),
    (fig2_path, 'Figure 2 - Institutional Logics & Accountability Patterns'),
    (fig3_path, 'Figure 3 - Ghost Node Exclusion Matrix'),
]:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    slide.shapes.add_picture(fig_path, margin, margin, img_width, img_height)
    print(f'  Added: {title}')

output_path = r'c:\Users\mount\.gemini\antigravity\scratch\Accountability_Architecture_Figures.pptx'
prs.save(output_path)
print(f'[OK] Saved: {output_path}')
