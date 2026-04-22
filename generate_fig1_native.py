"""
Figure 1: Directed Accountability Flows — Native PowerPoint shapes & text.
All elements are fully editable in PowerPoint.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ──────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────

def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

def add_rounded_rect(slide, left, top, width, height, fill_rgb, border_rgb, radius_pct=10, border_width=Pt(1.5)):
    """Add a rounded rectangle with fill and border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = border_width
    # Set corner radius - adjust avLst
    sp = shape._element
    prstGeom = sp.find('.//a:prstGeom', NSMAP)
    if prstGeom is not None:
        avLst = prstGeom.find('a:avLst', NSMAP)
        if avLst is None:
            avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for child in list(avLst):
                avLst.remove(child)
        gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj')
        gd.set('fmla', f'val {radius_pct * 1000}')
    # Remove text
    shape.text_frame.clear()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=10, bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if color:
        p.font.color.rgb = color
    return txBox

def add_node_box(slide, left, top, width, height, name, subtitle, fill_rgb=hex_to_rgb('#FFFFFF'), border_rgb=hex_to_rgb('#D0D5DD')):
    """Add a node box with name and subtitle text."""
    shape = add_rounded_rect(slide, left, top, width, height, fill_rgb, border_rgb, radius_pct=8, border_width=Pt(1.2))

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.05)
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    p1 = tf.paragraphs[0]
    p1.text = name
    p1.font.size = Pt(10)
    p1.font.bold = True
    p1.font.name = 'Calibri'
    p1.font.color.rgb = hex_to_rgb('#1A1A2E')
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(7.5)
    p2.font.bold = False
    p2.font.name = 'Calibri'
    p2.font.color.rgb = hex_to_rgb('#6B7280')
    p2.alignment = PP_ALIGN.CENTER

    return shape

def add_arrow_line(slide, x1, y1, x2, y2, color_hex, width_pt=2.0, dashed=False, head_end=True, head_start=False):
    """Add a line with arrowheads via direct XML manipulation."""
    # Create a connector (straight line)
    cxnSp = slide.shapes._spTree.makeelement(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cxnSp', {})

    # nvCxnSpPr
    nvCxnSpPr = etree.SubElement(cxnSp, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvCxnSpPr')
    cNvPr = etree.SubElement(nvCxnSpPr, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    cNvPr.set('id', str(100 + len(slide.shapes)))
    cNvPr.set('name', f'Arrow {len(slide.shapes)}')
    cNvCxnSpPr = etree.SubElement(nvCxnSpPr, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvCxnSpPr')
    nvPr = etree.SubElement(nvCxnSpPr, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')

    # spPr
    spPr = etree.SubElement(cxnSp, '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')

    # xfrm
    xfrm = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
    # Handle flipping for lines going right-to-left or bottom-to-top
    flip_h = '1' if x2 < x1 else '0'
    flip_v = '1' if y2 < y1 else '0'
    if flip_h == '1':
        xfrm.set('flipH', '1')
    if flip_v == '1':
        xfrm.set('flipV', '1')

    min_x = min(x1, x2)
    min_y = min(y1, y2)
    w = abs(x2 - x1)
    h = abs(y2 - y1)

    off = etree.SubElement(xfrm, '{http://schemas.openxmlformats.org/drawingml/2006/main}off')
    off.set('x', str(int(min_x)))
    off.set('y', str(int(min_y)))
    ext = etree.SubElement(xfrm, '{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
    ext.set('cx', str(int(w)))
    ext.set('cy', str(int(h)))

    # prstGeom = line
    prstGeom = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    prstGeom.set('prst', 'line')
    etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')

    # ln
    color = color_hex.lstrip('#')
    r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)

    ln = etree.SubElement(spPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    ln.set('w', str(int(width_pt * 12700)))

    solidFill = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    srgb = etree.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    srgb.set('val', f'{r:02X}{g:02X}{b:02X}')

    if dashed:
        prstDash = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
        prstDash.set('val', 'dash')

    if head_end:
        tailEnd = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', 'med')
        tailEnd.set('len', 'med')

    if head_start:
        headEnd = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
        headEnd.set('type', 'triangle')
        headEnd.set('w', 'med')
        headEnd.set('len', 'med')

    slide.shapes._spTree.append(cxnSp)

def add_legend_line(slide, x, y, width, color_hex, width_pt=2.5, dashed=False):
    """Add a horizontal line for the legend."""
    x1 = int(x)
    y1 = int(y)
    x2 = int(x + width)
    y2 = int(y)
    add_arrow_line(slide, x1, y1, x2, y2, color_hex, width_pt=width_pt, dashed=dashed, head_end=False)


# ──────────────────────────────────────────────────────────────
# LAYOUT CONSTANTS
# ──────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
BG_WHITE      = hex_to_rgb('#FFFFFF')
TIER_TOP_BG   = hex_to_rgb('#FEF3E7')
TIER_TOP_BD   = hex_to_rgb('#F0C8A0')
TIER_MID_BG   = hex_to_rgb('#EEF0FC')
TIER_MID_BD   = hex_to_rgb('#C0C4E8')
TIER_BOT_BG   = hex_to_rgb('#FDECEC')
TIER_BOT_BD   = hex_to_rgb('#E8B0B0')
NODE_BG       = hex_to_rgb('#FFFFFF')
NODE_BD       = hex_to_rgb('#D0D5DD')
ARROW_UP      = '#D95E0E'
ARROW_DOWN    = '#2980B9'
ARROW_ABSENT  = '#C0392B'
TEXT_DARK      = hex_to_rgb('#1A1A2E')
TEXT_DIM       = hex_to_rgb('#5A5F72')
TEXT_MUTED     = hex_to_rgb('#8A8FA2')
ACCENT_RED     = hex_to_rgb('#C0392B')

# Layout dimensions
MARGIN_X = Inches(1.2)
BOX_W = Inches(2.8)
BOX_H = Inches(0.65)
GAP_X = Inches(0.5)
TOTAL_ROW_W = 3 * BOX_W + 2 * GAP_X  # 9.4"
ROW_START_X = (SLIDE_W - TOTAL_ROW_W) // 2

TIER_W = TOTAL_ROW_W + Inches(0.6)
TIER_H = Inches(1.15)
TIER_X = ROW_START_X - Inches(0.3)

# Box center X positions
BOX_CX = [
    ROW_START_X + BOX_W // 2,
    ROW_START_X + BOX_W + GAP_X + BOX_W // 2,
    ROW_START_X + 2 * (BOX_W + GAP_X) + BOX_W // 2,
]

# ──────────────────────────────────────────────────────────────
# BUILD SLIDE 1
# ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = BG_WHITE

# ─── Title ───
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.4),
            'Figure 1. Directed Accountability Flows', font_size=18, bold=True, color=TEXT_DARK)

add_textbox(slide, Inches(0.8), Inches(0.55), Inches(11), Inches(0.35),
            '"Accountability flows through compliance architecture rather than enforceable obligations to those who have been impacted."',
            font_size=9, bold=False, color=TEXT_DIM, font_name='Calibri')
add_textbox(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.25),
            '— cf. Shaikh et al. (2023), relational spaces of algorithmic accountability',
            font_size=7.5, bold=False, color=TEXT_MUTED)

# ─── TIER 1: Oversight ───
T1_Y = Inches(1.15)
add_rounded_rect(slide, TIER_X, T1_Y, TIER_W, TIER_H, TIER_TOP_BG, TIER_TOP_BD, radius_pct=6)

add_textbox(slide, TIER_X + Inches(0.15), T1_Y + Inches(0.06), Inches(4), Inches(0.2),
            'OVERSIGHT & REGULATORY BODIES', font_size=7, bold=True, color=hex_to_rgb('#D95E0E'))

BOX_Y1 = T1_Y + Inches(0.35)
n_ra = add_node_box(slide, ROW_START_X, BOX_Y1, BOX_W, BOX_H, 'Regulatory Agencies', 'Standard-setting, enforcement')
n_cb = add_node_box(slide, ROW_START_X + BOX_W + GAP_X, BOX_Y1, BOX_W, BOX_H, 'Certification Bodies', 'Conformity assessment, audit')
n_nb = add_node_box(slide, ROW_START_X + 2*(BOX_W + GAP_X), BOX_Y1, BOX_W, BOX_H, 'Notified Bodies', 'Third-party verification')

# ─── TIER 2: Deployers ───
T2_Y = Inches(3.2)
add_rounded_rect(slide, TIER_X, T2_Y, TIER_W, TIER_H, TIER_MID_BG, TIER_MID_BD, radius_pct=6)

add_textbox(slide, TIER_X + Inches(0.15), T2_Y + Inches(0.06), Inches(4), Inches(0.2),
            'DEPLOYERS & OPERATORS', font_size=7, bold=True, color=hex_to_rgb('#6366F1'))

BOX_Y2 = T2_Y + Inches(0.35)
n_d  = add_node_box(slide, ROW_START_X, BOX_Y2, BOX_W, BOX_H, 'Deployers', 'Impact assessments, risk classification')
n_o  = add_node_box(slide, ROW_START_X + BOX_W + GAP_X, BOX_Y2, BOX_W, BOX_H, 'Operators', 'Documentation, monitoring')
n_pd = add_node_box(slide, ROW_START_X + 2*(BOX_W + GAP_X), BOX_Y2, BOX_W, BOX_H, 'Providers / Developers', 'Technical compliance, testing')

# ─── TIER 3: Impacted ───
T3_Y = Inches(5.25)
add_rounded_rect(slide, TIER_X, T3_Y, TIER_W, TIER_H, TIER_BOT_BG, TIER_BOT_BD, radius_pct=6)

add_textbox(slide, TIER_X + Inches(0.15), T3_Y + Inches(0.06), Inches(6), Inches(0.2),
            'IMPACTED COMMUNITIES & AFFECTED POPULATIONS', font_size=7, bold=True, color=hex_to_rgb('#C0392B'))

BOX_Y3 = T3_Y + Inches(0.35)
n_ic  = add_node_box(slide, ROW_START_X, BOX_Y3, BOX_W, BOX_H, 'Impacted Individuals', 'Subjects of algorithmic decisions')
n_cso = add_node_box(slide, ROW_START_X + BOX_W + GAP_X, BOX_Y3, BOX_W, BOX_H, 'Civil Society Orgs', 'Advocacy, representation')
n_mc  = add_node_box(slide, ROW_START_X + 2*(BOX_W + GAP_X), BOX_Y3, BOX_W, BOX_H, 'Marginalized Communities', 'Disproportionately affected groups')

# ──────────────────────────────────────────────────────────────
# ARROWS: Upward compliance flows (Mid → Top)
# ──────────────────────────────────────────────────────────────

# Helper: Get center-x of box n (0, 1, 2) with offset
def box_top(tier_box_y, offset=0):
    return int(tier_box_y + offset)

def box_bot(tier_box_y, box_h, offset=0):
    return int(tier_box_y + box_h + offset)

def box_cx(idx, offset_px=0):
    return int(BOX_CX[idx] + offset_px)

# Arrow endpoints
top_row_bot = box_bot(BOX_Y1, BOX_H)   # bottom of top tier boxes
mid_row_top = box_top(BOX_Y2)           # top of mid tier boxes
mid_row_bot = box_bot(BOX_Y2, BOX_H)   # bottom of mid tier boxes
bot_row_top = box_top(BOX_Y3)           # top of bot tier boxes

# 7 upward arrows (orange): from mid tier → top tier
offsets = [
    # (from_box_idx, from_x_offset, to_box_idx, to_x_offset, width)
    (0, -Inches(0.3), 0, 0,            2.5),  # D → RA (center)
    (0,  Inches(0.15), 1, -Inches(0.3), 1.8),  # D → CB
    (0,  Inches(0.5), 2, -Inches(0.4), 1.5),  # D → NB
    (1, -Inches(0.3), 0,  Inches(0.35), 1.8),  # O → RA
    (1,  0,           1, 0,            2.5),  # O → CB (center)
    (2, -Inches(0.3), 1,  Inches(0.3), 1.8),  # P → CB
    (2,  Inches(0.1), 2, 0,            2.5),  # P → NB (center)
]

for from_idx, from_off, to_idx, to_off, w in offsets:
    add_arrow_line(slide,
        box_cx(from_idx, from_off), mid_row_top,
        box_cx(to_idx, to_off), top_row_bot,
        ARROW_UP, width_pt=w)

# 2 downward arrows (blue, dashed): from top tier → mid tier
add_arrow_line(slide,
    box_cx(0, -Inches(0.6)), top_row_bot,
    box_cx(0, -Inches(0.65)), mid_row_top,
    ARROW_DOWN, width_pt=1.5, dashed=True)

add_arrow_line(slide,
    box_cx(1, Inches(0.5)), top_row_bot,
    box_cx(1, Inches(0.45)), mid_row_top,
    ARROW_DOWN, width_pt=1.5, dashed=True)

# 3 absent arrows (red, dashed): between mid ↔ bot
# IC → D (upward, no channel)
add_arrow_line(slide,
    box_cx(0, -Inches(0.1)), bot_row_top,
    box_cx(0, -Inches(0.15)), mid_row_bot,
    ARROW_ABSENT, width_pt=1.5, dashed=True)

# CSO → O (upward, no standing)
add_arrow_line(slide,
    box_cx(1, -Inches(0.1)), bot_row_top,
    box_cx(1, -Inches(0.15)), mid_row_bot,
    ARROW_ABSENT, width_pt=1.5, dashed=True)

# D → IC (downward, no obligation)
add_arrow_line(slide,
    box_cx(0, Inches(0.3)), mid_row_bot,
    box_cx(0, Inches(0.25)), bot_row_top,
    ARROW_ABSENT, width_pt=1.5, dashed=True)

# ✗ markers between mid and bot
mid_bot_gap_y = (mid_row_bot + bot_row_top) // 2
add_textbox(slide, BOX_CX[0] - Inches(0.25), mid_bot_gap_y - Inches(0.08), Inches(0.3), Inches(0.2),
            '✗', font_size=10, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)
add_textbox(slide, BOX_CX[0] + Inches(0.1), mid_bot_gap_y - Inches(0.08), Inches(0.3), Inches(0.2),
            '✗', font_size=10, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)
add_textbox(slide, BOX_CX[1] - Inches(0.15), mid_bot_gap_y - Inches(0.08), Inches(0.3), Inches(0.2),
            '✗', font_size=10, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────
# LEGEND
# ──────────────────────────────────────────────────────────────

LEG_Y = Inches(6.55)
leg_line_w = Inches(0.5)

# Upward
add_legend_line(slide, Inches(2.5), LEG_Y + Inches(0.08), leg_line_w, ARROW_UP, width_pt=2.5)
add_textbox(slide, Inches(3.1), LEG_Y - Inches(0.02), Inches(2.8), Inches(0.22),
            'Strong upward compliance flow (7 arrows)', font_size=8, color=TEXT_DIM)

# Downward
add_legend_line(slide, Inches(6.2), LEG_Y + Inches(0.08), leg_line_w, ARROW_DOWN, width_pt=1.5, dashed=True)
add_textbox(slide, Inches(6.8), LEG_Y - Inches(0.02), Inches(2.5), Inches(0.22),
            'Weak downward obligation (2 arrows)', font_size=8, color=TEXT_DIM)

# Absent
add_legend_line(slide, Inches(9.5), LEG_Y + Inches(0.08), leg_line_w, ARROW_ABSENT, width_pt=1.5, dashed=True)
add_textbox(slide, Inches(10.1), LEG_Y - Inches(0.02), Inches(2.5), Inches(0.22),
            'Absent / no formal pathway (3 arrows)', font_size=8, color=TEXT_DIM)

# ──────────────────────────────────────────────────────────────
# CALLOUT
# ──────────────────────────────────────────────────────────────

CALL_Y = Inches(6.85)
callout_bg = add_rounded_rect(slide, Inches(1.0), CALL_Y, Inches(11.3), Inches(0.5),
                               hex_to_rgb('#FEF2F2'), hex_to_rgb('#FECACA'), radius_pct=5, border_width=Pt(1))

# Red left bar accent (narrow rectangle)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), CALL_Y, Inches(0.06), Inches(0.5))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_RED
bar.line.fill.background()

# Callout text
tf = callout_bg.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_right = Inches(0.15)
tf.margin_top = Inches(0.06)

p1 = tf.paragraphs[0]
run_bold = p1.add_run()
run_bold.text = 'Structural observation: '
run_bold.font.size = Pt(8)
run_bold.font.bold = True
run_bold.font.color.rgb = ACCENT_RED
run_bold.font.name = 'Calibri'

run_text = p1.add_run()
run_text.text = ('The density of upward arrows (compliance reporting, conformity assessments, audit documentation) '
                 'vastly exceeds downward flows. Impacted communities have ')
run_text.font.size = Pt(8)
run_text.font.color.rgb = TEXT_DIM
run_text.font.name = 'Calibri'

run_bold2 = p1.add_run()
run_bold2.text = 'no formal channels'
run_bold2.font.size = Pt(8)
run_bold2.font.bold = True
run_bold2.font.color.rgb = ACCENT_RED
run_bold2.font.name = 'Calibri'

run_end = p1.add_run()
run_end.text = ' to influence outcomes or hold deployers accountable — accountability is relational to regulators, not to those harmed.'
run_end.font.size = Pt(8)
run_end.font.color.rgb = TEXT_DIM
run_end.font.name = 'Calibri'

# ──────────────────────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────────────────────

output_path = r'c:\Users\mount\.gemini\antigravity\scratch\Figure1_Native_Shapes.pptx'
prs.save(output_path)
print(f'[OK] Saved: {output_path}')
print('All shapes, arrows, and text are native PowerPoint objects — fully editable.')
