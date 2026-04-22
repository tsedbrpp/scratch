"""
Worker Exclusion — Rhetoric-Architecture Gap Chart
Five jurisdictions showing the disconnect between rhetorical invocation
of workers and their institutional role in AI governance architectures.

All elements are native PowerPoint shapes — fully editable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# ──────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, left, top, width, height, fill, border, radius=6, bw=Pt(1)):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = bw
    sp = sh._element
    pg = sp.find('.//a:prstGeom', NSMAP)
    if pg is not None:
        av = pg.find('a:avLst', NSMAP)
        if av is None:
            av = etree.SubElement(pg, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for c in list(av): av.remove(c)
        gd = etree.SubElement(av, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj'); gd.set('fmla', f'val {radius * 1000}')
    sh.text_frame.clear()
    return sh

def add_text(slide, left, top, width, height, text, size=10, bold=False, color=None,
             align=PP_ALIGN.LEFT, font='Calibri', italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.bold = bold; p.font.name = font; p.alignment = align
    p.font.italic = italic
    if color: p.font.color.rgb = color
    return tb

def rich_cell(slide, x, y, w, h, bg, bd, lines, bw=Pt(1), ml=0.1, mt=0.08, radius=5, anchor='top'):
    card = add_rect(slide, x, y, w, h, bg, bd, radius=radius, bw=bw)
    tf = card.text_frame
    tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Inches(ml); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(0.04)
    if anchor == 'middle':
        bodyPr = tf._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
    first = True
    for ln in lines:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln.get('t', '')
        r.font.size = Pt(ln.get('s', 10))
        r.font.bold = ln.get('b', False)
        r.font.italic = ln.get('i', False)
        r.font.name = 'Calibri'
        if 'c' in ln: r.font.color.rgb = hex_rgb(ln['c'])
        p.space_before = Pt(ln.get('sb', 1))
        p.space_after = Pt(ln.get('sa', 1))
        p.alignment = ln.get('a', PP_ALIGN.LEFT)
    return card

def add_bar(slide, x, y, w, h, fill_hex, border_hex=None):
    """Add a simple filled rectangle (no rounding) for rhetoric volume bars."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = hex_rgb(fill_hex)
    if border_hex:
        sh.line.color.rgb = hex_rgb(border_hex); sh.line.width = Pt(0.5)
    else:
        sh.line.fill.background()
    sh.text_frame.clear()
    return sh

def add_arrow(slide, x1, y1, x2, y2, color_hex, w=2.0, dashed=False, head=True):
    cxn = slide.shapes._spTree.makeelement(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cxnSp', {})
    nv = etree.SubElement(cxn, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvCxnSpPr')
    cp = etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    cp.set('id', str(200+len(slide.shapes))); cp.set('name', f'Arr{len(slide.shapes)}')
    etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvCxnSpPr')
    etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
    sp = etree.SubElement(cxn, '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    xf = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
    if x2 < x1: xf.set('flipH', '1')
    if y2 < y1: xf.set('flipV', '1')
    off = etree.SubElement(xf, '{http://schemas.openxmlformats.org/drawingml/2006/main}off')
    off.set('x', str(int(min(x1,x2)))); off.set('y', str(int(min(y1,y2))))
    ext = etree.SubElement(xf, '{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
    ext.set('cx', str(int(abs(x2-x1)))); ext.set('cy', str(int(abs(y2-y1))))
    pg = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    pg.set('prst', 'line')
    etree.SubElement(pg, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
    c = color_hex.lstrip('#')
    ln = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    ln.set('w', str(int(w * 12700)))
    sf = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    sr = etree.SubElement(sf, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    sr.set('val', c.upper())
    if dashed:
        pd = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
        pd.set('val', 'dash')
    if head:
        te = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
        te.set('type', 'triangle'); te.set('w', 'lg'); te.set('len', 'lg')
    slide.shapes._spTree.append(cxn)


# ──────────────────────────────────────────────────────────────
# DATA
# ──────────────────────────────────────────────────────────────

JURISDICTIONS = [
    {
        'name': 'European Union',
        'color': '#6366F1',
        'light': '#EEF2FF',
        'border': '#A5B4FC',
        'rhetoric_level': 0.55,   # moderate — workers mentioned in recitals, workplace AI flagged
        'rhetoric_lines': [
            {'t': 'Recital 57 references AI in', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'the workplace; "employment"', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'listed in Annex III high-risk', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'categories; Art. 26 deployer', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'obligations mention workers', 's': 9, 'b': False, 'c': '#78350F'},
        ],
        'institutional_lines': [
            {'t': '✗  No role in conformity', 's': 9, 'b': True, 'c': '#991B1B'},
            {'t': '   assessment process', 's': 9, 'b': False, 'c': '#991B1B'},
            {'t': '✗  Excluded from compliance', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '   assessment design', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '✗  No seat on AI Board', 's': 9, 'b': False, 'c': '#7F1D1D'},
        ],
        'ghost_score': '87',
    },
    {
        'name': 'Colorado',
        'color': '#16A34A',
        'light': '#F0FDF4',
        'border': '#86EFAC',
        'rhetoric_level': 0.40,
        'rhetoric_lines': [
            {'t': '"Consequential decisions"', 's': 9, 'b': True, 'c': '#78350F'},
            {'t': 'explicitly includes', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'employment decisions;', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'deployers must assess impact', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'on affected individuals', 's': 9, 'b': False, 'c': '#78350F'},
        ],
        'institutional_lines': [
            {'t': '✗  No worker voice in', 's': 9, 'b': True, 'c': '#991B1B'},
            {'t': '   impact assessments', 's': 9, 'b': False, 'c': '#991B1B'},
            {'t': '✗  AG sole enforcer — no', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '   private right of action', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '✗  Self-testing regime', 's': 9, 'b': False, 'c': '#7F1D1D'},
        ],
        'ghost_score': '—',
    },
    {
        'name': 'India',
        'color': '#D97706',
        'light': '#FFFBEB',
        'border': '#FCD34D',
        'rhetoric_level': 0.25,
        'rhetoric_lines': [
            {'t': 'NITI Aayog principles', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'reference "inclusive AI"', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'and labor disruption;', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'developmental framing', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'subsumes worker concerns', 's': 9, 'b': False, 'c': '#78350F'},
        ],
        'institutional_lines': [
            {'t': '✗  No labor channel in', 's': 9, 'b': True, 'c': '#991B1B'},
            {'t': '   governance architecture', 's': 9, 'b': False, 'c': '#991B1B'},
            {'t': '✗  Enterprise-only design;', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '   deployment without labor', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '   impact analysis', 's': 9, 'b': False, 'c': '#7F1D1D'},
        ],
        'ghost_score': '—',
    },
    {
        'name': 'Brazil',
        'color': '#0D9488',
        'light': '#F0FDFA',
        'border': '#5EEAD4',
        'rhetoric_level': 0.50,
        'rhetoric_lines': [
            {'t': 'PL 2338/2023 invokes', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'workers\' rights via LGPD', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'framework; formal right to', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'explanation and human review', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'mentioned in text', 's': 9, 'b': False, 'c': '#78350F'},
        ],
        'institutional_lines': [
            {'t': '✗  Principle-level only;', 's': 9, 'b': True, 'c': '#991B1B'},
            {'t': '   no institutional seat', 's': 9, 'b': False, 'c': '#991B1B'},
            {'t': '✗  ANPD coordination has', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '   no labor representative', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '✗  Compliance-oriented', 's': 9, 'b': False, 'c': '#7F1D1D'},
        ],
        'ghost_score': '80',
    },
    {
        'name': 'United States',
        'color': '#6B7280',
        'light': '#F9FAFB',
        'border': '#9CA3AF',
        'rhetoric_level': 0.95,   # extreme — most rhetorical commitment of all five
        'rhetoric_lines': [
            {'t': '"Worker-first AI agenda"', 's': 9.5, 'b': True, 'c': '#78350F'},
            {'t': '"Workers must benefit from', 's': 9, 'b': True, 'c': '#92400E', 'i': True},
            {'t': ' AI-driven growth"', 's': 9, 'b': True, 'c': '#92400E', 'i': True},
            {'t': 'Executive Orders center', 's': 9, 'b': False, 'c': '#78350F'},
            {'t': 'worker welfare prominently', 's': 9, 'b': False, 'c': '#78350F'},
        ],
        'institutional_lines': [
            {'t': '✗  No role in compliance', 's': 9, 'b': True, 'c': '#991B1B'},
            {'t': '   design or evaluation', 's': 9, 'b': False, 'c': '#991B1B'},
            {'t': '✗  No seat on governance', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '   boards (CAISI, CAIOC)', 's': 9, 'b': False, 'c': '#7F1D1D'},
            {'t': '✗  No institutional channel', 's': 9, 'b': False, 'c': '#7F1D1D'},
        ],
        'ghost_score': '—',
    },
]


# ──────────────────────────────────────────────────────────────
# LAYOUT
# ──────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

L_MARGIN = Inches(0.35)
R_MARGIN = Inches(0.35)
USABLE_W = SLIDE_W - L_MARGIN - R_MARGIN

# Column widths
JUR_NAME_W = Inches(1.4)         # Jurisdiction label column
BAR_W = Inches(0.22)             # Rhetoric volume bar
RHETORIC_W = Inches(3.8)         # Rhetoric text panel
GAP_W = Inches(0.12)             # Gap between panels
DIVIDER_W = Inches(0.06)         # Central divider
INSTIT_W = Inches(3.8)           # Institutional role panel
CROSS_W = Inches(0.7)            # Cross/exclusion indicator
SCORE_W = Inches(0.7)            # Ghost node score

# Vertical layout
TITLE_TOP = Inches(0.12)
HEADER_Y = Inches(0.88)
HEADER_H = Inches(0.38)
ROW_GAP = Inches(0.06)
ROW_H = Inches(1.05)
ROW1_TOP = HEADER_Y + HEADER_H + ROW_GAP

# Column X positions
JUR_X = L_MARGIN
BAR_X = JUR_X + JUR_NAME_W + Inches(0.06)
RHET_X = BAR_X + BAR_W + Inches(0.06)
DIVIDER_X = RHET_X + RHETORIC_W + GAP_W
INSTIT_X = DIVIDER_X + DIVIDER_W + GAP_W
SCORE_X = INSTIT_X + INSTIT_W + Inches(0.12)

# ──────────────────────────────────────────────────────────────
# BUILD
# ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = hex_rgb('#FFFFFF')

# ─── Title ───
add_text(slide, Inches(0.4), TITLE_TOP, Inches(12), Inches(0.4),
         'Worker Exclusion — Rhetoric vs. Institutional Architecture',
         size=20, bold=True, color=hex_rgb('#111827'))

add_text(slide, Inches(0.4), TITLE_TOP + Inches(0.38), Inches(12), Inches(0.35),
         'The recurring mechanism is not merely an omission but a structured exclusion from the architectures through which governance becomes actionable.',
         size=10.5, bold=False, color=hex_rgb('#6B7280'), italic=True)

# ─── Column Headers ───
# Jurisdiction
rich_cell(slide, JUR_X, HEADER_Y, JUR_NAME_W, HEADER_H,
          hex_rgb('#1E293B'), hex_rgb('#0F172A'),
          [{'t': 'JURISDICTION', 's': 8, 'b': True, 'c': '#F8FAFC', 'a': PP_ALIGN.CENTER}],
          bw=Pt(0), mt=0.06, anchor='middle')

# Rhetoric header
rhetoric_hdr_w = BAR_W + Inches(0.06) + RHETORIC_W
rich_cell(slide, BAR_X, HEADER_Y, rhetoric_hdr_w, HEADER_H,
          hex_rgb('#FFFBEB'), hex_rgb('#FDE68A'),
          [{'t': 'RHETORICAL PRESENCE', 's': 8, 'b': True, 'c': '#92400E', 'a': PP_ALIGN.CENTER},
           {'t': 'What the policy text says about workers', 's': 7, 'b': False, 'c': '#B45309',
            'a': PP_ALIGN.CENTER, 'sb': 0}],
          bw=Pt(1), mt=0.04, anchor='middle')

# Institutional header
rich_cell(slide, INSTIT_X, HEADER_Y, INSTIT_W, HEADER_H,
          hex_rgb('#FEF2F2'), hex_rgb('#FECACA'),
          [{'t': 'INSTITUTIONAL ROLE', 's': 8, 'b': True, 'c': '#991B1B', 'a': PP_ALIGN.CENTER},
           {'t': 'What the governance architecture provides', 's': 7, 'b': False, 'c': '#B91C1C',
            'a': PP_ALIGN.CENTER, 'sb': 0}],
          bw=Pt(1), mt=0.04, anchor='middle')

# Ghost score header
rich_cell(slide, SCORE_X, HEADER_Y, SCORE_W, HEADER_H,
          hex_rgb('#1E293B'), hex_rgb('#0F172A'),
          [{'t': 'GHOST\nSCORE', 's': 7, 'b': True, 'c': '#F8FAFC', 'a': PP_ALIGN.CENTER}],
          bw=Pt(0), mt=0.04, anchor='middle')

# ─── Central Divider Line (vertical) ───
divider_top = int(HEADER_Y)
divider_bot = int(ROW1_TOP + 5 * ROW_H + 4 * ROW_GAP)
div_cx = int(DIVIDER_X + DIVIDER_W // 2)
add_arrow(slide, div_cx, divider_top, div_cx, divider_bot, '#D1D5DB', w=1.5, dashed=True, head=False)

# ─── Data Rows ───
for ri, jur in enumerate(JURISDICTIONS):
    ry = ROW1_TOP + ri * (ROW_H + ROW_GAP)
    jcolor = jur['color']
    jlight = jur['light']
    jborder = jur['border']
    is_us = (jur['name'] == 'United States')

    # ── Row background strip ──
    if is_us:
        # U.S. row gets special emphasis
        row_bg = add_rect(slide, L_MARGIN - Inches(0.08), ry - Inches(0.04),
                          USABLE_W + Inches(0.16), ROW_H + Inches(0.08),
                          hex_rgb('#FFF7ED'), hex_rgb('#FED7AA'), radius=6, bw=Pt(1.5))

    # ── Jurisdiction name cell ──
    jur_cell_bg = '#FFF7ED' if is_us else jlight
    jur_cell_bd = '#FED7AA' if is_us else jborder
    name_lines = [
        {'t': jur['name'], 's': 11, 'b': True, 'c': jcolor, 'a': PP_ALIGN.CENTER},
    ]
    rich_cell(slide, JUR_X, ry, JUR_NAME_W, ROW_H,
              hex_rgb(jur_cell_bg), hex_rgb(jur_cell_bd),
              name_lines, bw=Pt(1.5), mt=0.06, anchor='middle')

    # Add jurisdiction dot inside name cell
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  JUR_X + JUR_NAME_W // 2 - Inches(0.05),
                                  ry + Inches(0.08),
                                  Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = hex_rgb(jcolor)
    dot.line.fill.background()

    # ── Rhetoric volume bar ──
    max_bar_h = ROW_H - Inches(0.12)
    bar_h = int(max_bar_h * jur['rhetoric_level'])
    bar_y = ry + ROW_H - Inches(0.06) - bar_h  # Grow from bottom

    # Gradient-like effect: amber intensity proportional to rhetoric level
    if jur['rhetoric_level'] >= 0.8:
        bar_color = '#F59E0B'   # Strong amber
        bar_border = '#D97706'
    elif jur['rhetoric_level'] >= 0.5:
        bar_color = '#FCD34D'   # Medium amber
        bar_border = '#FBBF24'
    else:
        bar_color = '#FDE68A'   # Light amber
        bar_border = '#FCD34D'

    add_bar(slide, BAR_X + Inches(0.03), bar_y, BAR_W - Inches(0.06), bar_h,
            bar_color, bar_border)

    # ── Rhetoric text panel ──
    rhet_bg = '#FFFBEB' if not is_us else '#FEF3C7'
    rhet_bd = '#FDE68A' if not is_us else '#FBBF24'
    rich_cell(slide, RHET_X, ry, RHETORIC_W, ROW_H,
              hex_rgb(rhet_bg), hex_rgb(rhet_bd),
              jur['rhetoric_lines'],
              bw=Pt(1), mt=0.08, ml=0.12, anchor='middle')

    # ── Institutional role panel ──
    inst_bg = '#FEF2F2'
    inst_bd = '#FECACA'
    rich_cell(slide, INSTIT_X, ry, INSTIT_W, ROW_H,
              hex_rgb(inst_bg), hex_rgb(inst_bd),
              jur['institutional_lines'],
              bw=Pt(1), mt=0.08, ml=0.12, anchor='middle')

    # ── Ghost node score ──
    score = jur['ghost_score']
    if score != '—':
        score_lines = [
            {'t': score, 's': 20, 'b': True, 'c': '#991B1B', 'a': PP_ALIGN.CENTER},
            {'t': '/100', 's': 8, 'b': False, 'c': '#9CA3AF', 'a': PP_ALIGN.CENTER, 'sb': 0},
        ]
        score_bg = '#FEF2F2'
        score_bd = '#FECACA'
    else:
        score_lines = [
            {'t': '—', 's': 18, 'b': False, 'c': '#D1D5DB', 'a': PP_ALIGN.CENTER},
        ]
        score_bg = '#FAFBFC'
        score_bd = '#E5E7EB'

    rich_cell(slide, SCORE_X, ry, SCORE_W, ROW_H,
              hex_rgb(score_bg), hex_rgb(score_bd),
              score_lines,
              bw=Pt(0.75), mt=0.06, anchor='middle')

# ─── U.S. Case Annotation Arrow ───
us_row_y = ROW1_TOP + 4 * (ROW_H + ROW_GAP)
annot_x = SCORE_X + SCORE_W + Inches(0.1)
annot_y = us_row_y + ROW_H // 2

# Arrow pointing to the U.S. row
add_arrow(slide, int(annot_x + Inches(1.6)), int(annot_y),
          int(annot_x + Inches(0.08)), int(annot_y),
          '#D97706', w=2.0)

add_text(slide, annot_x + Inches(0.2), annot_y - Inches(0.35), Inches(1.9), Inches(0.7),
         'Maximum rhetoric,\nzero institutional\nrole — sharpest case',
         size=8, bold=True, color=hex_rgb('#B45309'), align=PP_ALIGN.LEFT, italic=True)

# ──────────────────────────────────────────────────────────────
# CONVERGENCE CALLOUT
# ──────────────────────────────────────────────────────────────

call_y = ROW1_TOP + 5 * (ROW_H + ROW_GAP) + Inches(0.08)
call_w = SCORE_X + SCORE_W - L_MARGIN
call_h = Inches(0.5)

call_bg = add_rect(slide, L_MARGIN, call_y, call_w, call_h,
                    hex_rgb('#FEF2F2'), hex_rgb('#FECACA'), radius=5, bw=Pt(1.2))

# Red accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L_MARGIN, call_y, Inches(0.06), call_h)
bar.fill.solid(); bar.fill.fore_color.rgb = hex_rgb('#991B1B')
bar.line.fill.background()

tf = call_bg.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.15)
tf.margin_top = Inches(0.06)

p = tf.paragraphs[0]

r1 = p.add_run()
r1.text = 'Structural observation: '
r1.font.size = Pt(8.5); r1.font.bold = True
r1.font.color.rgb = hex_rgb('#991B1B'); r1.font.name = 'Calibri'

r2 = p.add_run()
r2.text = ('All five jurisdictions govern domains in which AI affects employment, labor process, '
           'or workplace decision-making. Yet none creates an institutional role for worker organizations '
           'in compliance design, evaluation, or governance boards. ')
r2.font.size = Pt(8.5); r2.font.color.rgb = hex_rgb('#7F1D1D'); r2.font.name = 'Calibri'

r3 = p.add_run()
r3.text = 'The exclusion is structured, not incidental.'
r3.font.size = Pt(8.5); r3.font.bold = True; r3.font.italic = True
r3.font.color.rgb = hex_rgb('#991B1B'); r3.font.name = 'Calibri'

# ──────────────────────────────────────────────────────────────
# LEGEND
# ──────────────────────────────────────────────────────────────

leg_y = call_y + call_h + Inches(0.14)

# Rhetoric bar legend
add_bar(slide, L_MARGIN + Inches(0.3), leg_y + Inches(0.03), Inches(0.5), Inches(0.12),
        '#FCD34D', '#FBBF24')
add_text(slide, L_MARGIN + Inches(0.9), leg_y - Inches(0.02), Inches(3.5), Inches(0.2),
         'Rhetorical commitment level (bar height = relative intensity)',
         size=8, color=hex_rgb('#78350F'))

# Institutional absence
add_text(slide, L_MARGIN + Inches(4.6), leg_y - Inches(0.02), Inches(0.3), Inches(0.2),
         '✗', size=10, bold=True, color=hex_rgb('#991B1B'), align=PP_ALIGN.CENTER)
add_text(slide, L_MARGIN + Inches(4.95), leg_y - Inches(0.02), Inches(3), Inches(0.2),
         'No institutional role in governance function',
         size=8, color=hex_rgb('#991B1B'))

# Ghost score
add_text(slide, L_MARGIN + Inches(8.3), leg_y - Inches(0.02), Inches(4), Inches(0.2),
         'Ghost Score — absence strength from Ghost Node Detection Pipeline (where measured)',
         size=8, color=hex_rgb('#6B7280'))

# ──────────────────────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────────────────────

output = r'c:\Users\mount\.gemini\antigravity\scratch\Worker_Exclusion_Rhetoric_Gap.pptx'
prs.save(output)
print(f'[OK] Saved: {output}')
print('All shapes, arrows, and text are native PowerPoint objects — fully editable.')
