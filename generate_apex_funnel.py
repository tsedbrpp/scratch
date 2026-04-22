"""
Authority Concentration — Apex Funnel Diagram
Five parallel funnels showing convergent authority compression across
EU, Colorado, India, Brazil, and the U.S.

All elements are native PowerPoint shapes — fully editable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# ──────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, left, top, width, height, fill, border, radius=6, bw=Pt(1)):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = bw
    sp = sh._element
    pg = sp.find('.//a:prstGeom', NSMAP)
    if pg is not None:
        av = pg.find('a:avLst', NSMAP)
        if av is None:
            av = etree.SubElement(pg, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for c in list(av): av.remove(c)
        gd = etree.SubElement(av, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj'); gd.set('fmla', f'val {radius * 1000}')
    sh.text_frame.clear()
    return sh

def add_sharp_rect(slide, left, top, width, height, fill, border, bw=Pt(1)):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = bw
    sh.text_frame.clear()
    return sh

def add_text(slide, left, top, width, height, text, size=10, bold=False, color=None,
             align=PP_ALIGN.LEFT, font='Calibri', italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.bold = bold; p.font.name = font; p.alignment = align
    p.font.italic = italic
    if color: p.font.color.rgb = color
    return tb

def rich_text(slide, x, y, w, h, lines, ml=0.1, mt=0.06):
    """Multi-run styled textbox."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Inches(ml); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(0.04)
    first = True
    for ln in lines:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln.get('t', '')
        r.font.size = Pt(ln.get('s', 10))
        r.font.bold = ln.get('b', False)
        r.font.italic = ln.get('i', False)
        r.font.name = 'Calibri'
        if 'c' in ln: r.font.color.rgb = hex_rgb(ln['c'])
        p.space_before = Pt(ln.get('sb', 1))
        p.space_after = Pt(ln.get('sa', 1))
        p.alignment = ln.get('a', PP_ALIGN.LEFT)
    return tb

def rich_cell(slide, x, y, w, h, bg, bd, lines, bw=Pt(1), ml=0.1, mt=0.08, radius=5, anchor='top'):
    card = add_rect(slide, x, y, w, h, bg, bd, radius=radius, bw=bw)
    tf = card.text_frame
    tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Inches(ml); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(0.04)
    if anchor == 'middle':
        bodyPr = tf._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
    first = True
    for ln in lines:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln.get('t', '')
        r.font.size = Pt(ln.get('s', 10))
        r.font.bold = ln.get('b', False)
        r.font.italic = ln.get('i', False)
        r.font.name = 'Calibri'
        if 'c' in ln: r.font.color.rgb = hex_rgb(ln['c'])
        p.space_before = Pt(ln.get('sb', 1))
        p.space_after = Pt(ln.get('sa', 1))
        p.alignment = ln.get('a', PP_ALIGN.LEFT)
    return card

def add_arrow(slide, x1, y1, x2, y2, color_hex, w=2.0, dashed=False, head=True):
    cxn = slide.shapes._spTree.makeelement(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cxnSp', {})
    nv = etree.SubElement(cxn, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvCxnSpPr')
    cp = etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    cp.set('id', str(200+len(slide.shapes))); cp.set('name', f'Arr{len(slide.shapes)}')
    etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvCxnSpPr')
    etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
    sp = etree.SubElement(cxn, '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    xf = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
    if x2 < x1: xf.set('flipH', '1')
    if y2 < y1: xf.set('flipV', '1')
    off = etree.SubElement(xf, '{http://schemas.openxmlformats.org/drawingml/2006/main}off')
    off.set('x', str(int(min(x1,x2)))); off.set('y', str(int(min(y1,y2))))
    ext = etree.SubElement(xf, '{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
    ext.set('cx', str(int(abs(x2-x1)))); ext.set('cy', str(int(abs(y2-y1))))
    pg = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    pg.set('prst', 'line')
    etree.SubElement(pg, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
    c = color_hex.lstrip('#')
    ln = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    ln.set('w', str(int(w * 12700)))
    sf = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    sr = etree.SubElement(sf, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    sr.set('val', c.upper())
    if dashed:
        pd = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
        pd.set('val', 'dash')
    if head:
        te = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
        te.set('type', 'triangle'); te.set('w', 'lg'); te.set('len', 'lg')
    slide.shapes._spTree.append(cxn)

def add_trapezoid(slide, left, top, width, height, fill_rgb, border_rgb, bw=Pt(1)):
    """Add a trapezoid shape (wider at bottom, narrower at top)."""
    sh = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_rgb
    sh.line.color.rgb = border_rgb; sh.line.width = bw
    # Flip vertically by rotating 180 degrees to get funnel shape (narrow top, wide bottom)
    sh.rotation = 180.0
    sh.text_frame.clear()
    return sh

def add_triangle(slide, left, top, width, height, fill_rgb, border_rgb, bw=Pt(1)):
    """Add an isosceles triangle pointing upward."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill_rgb
    sh.line.color.rgb = border_rgb; sh.line.width = bw
    sh.text_frame.clear()
    return sh

def add_chevron_up(slide, cx, y, size, color_hex):
    """Add a small upward chevron arrow using a textbox."""
    add_text(slide, cx - Inches(0.12), y, Inches(0.24), Inches(0.2),
             '▲', size=9, bold=True, color=hex_rgb(color_hex), align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────────
# JURISDICTION DATA
# ──────────────────────────────────────────────────────────────

JURISDICTIONS = [
    {
        'name': 'European Union',
        'short': 'EU',
        'color': '#6366F1',        # indigo
        'color_light': '#EEF2FF',
        'color_mid': '#C7D2FE',
        'color_border': '#A5B4FC',
        'apex': 'European Commission\n+ AI Office',
        'mid': ['CEN / CENELEC', 'Conformity\nassessment'],
        'base': ['AI Act risk tiers', 'CE marking', 'Deployer\nobligations'],
        'arch_type': 'Regulatory\npyramid',
    },
    {
        'name': 'Colorado',
        'short': 'CO',
        'color': '#16A34A',        # green
        'color_light': '#F0FDF4',
        'color_mid': '#BBF7D0',
        'color_border': '#86EFAC',
        'apex': 'Attorney\nGeneral',
        'mid': ['Impact\nassessments', 'Duty of care'],
        'base': ['SB 21-169', 'Deployer\nduties', 'No private\nright of action'],
        'arch_type': 'Monocentric\nenforcement',
    },
    {
        'name': 'India',
        'short': 'IN',
        'color': '#D97706',        # amber
        'color_light': '#FFFBEB',
        'color_mid': '#FDE68A',
        'color_border': '#FCD34D',
        'apex': 'CARO + AIREC\n+ IndiaAI Mission',
        'mid': ['Ministerial\ncoordination', 'NITI Aayog\nprinciples'],
        'base': ['DPDP Act', 'Development\npriorities', 'Industry\nself-regulation'],
        'arch_type': 'Multi-hub\ncoordination',
    },
    {
        'name': 'Brazil',
        'short': 'BR',
        'color': '#0D9488',        # teal
        'color_light': '#F0FDFA',
        'color_mid': '#99F6E4',
        'color_border': '#5EEAD4',
        'apex': 'ANPD',
        'mid': ['Sectoral\nregulators', 'LGPD-derived\naudit routines'],
        'base': ['PL 2338/2023', 'Rights\nframework', 'Multi-agency\necosystem'],
        'arch_type': 'Centralized\ncoordinator',
    },
    {
        'name': 'United States',
        'short': 'US',
        'color': '#6B7280',        # slate/gray
        'color_light': '#F9FAFB',
        'color_mid': '#D1D5DB',
        'color_border': '#9CA3AF',
        'apex': 'OMB + CAISI\n+ CAIOC',
        'mid': ['NIST\nframeworks', 'Procurement\ninfrastructure'],
        'base': ['EO 14110', 'Agency\nimplementation', 'No standalone\nregulator'],
        'arch_type': 'Distributed\nconstellation',
    },
]

# ──────────────────────────────────────────────────────────────
# LAYOUT CONSTANTS
# ──────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NUM_COLS = 5
COL_GAP = Inches(0.16)
L_MARGIN = Inches(0.4)
R_MARGIN = Inches(0.4)
USABLE_W = SLIDE_W - L_MARGIN - R_MARGIN
COL_W = (USABLE_W - (NUM_COLS - 1) * COL_GAP) // NUM_COLS

# Vertical layout
TITLE_TOP = Inches(0.15)
HEADER_TOP = Inches(0.92)        # Jurisdiction name + architecture type
HEADER_H = Inches(0.55)

# Funnel layers (from top to bottom: apex, mid, base)
APEX_TOP = Inches(1.65)
APEX_H = Inches(0.9)
MID_TOP = Inches(3.05)
MID_H = Inches(1.15)
BASE_TOP = Inches(4.7)
BASE_H = Inches(1.2)

# Funnel visual element positions
FUNNEL_TOP = Inches(1.55)
FUNNEL_BOT = Inches(5.95)

# ──────────────────────────────────────────────────────────────
# BUILD
# ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = hex_rgb('#FFFFFF')

# ─── Title ───
add_text(slide, Inches(0.5), TITLE_TOP, Inches(12), Inches(0.4),
         'Authority Concentration — Apex Funnel Architecture',
         size=20, bold=True, color=hex_rgb('#111827'))

add_text(slide, Inches(0.5), TITLE_TOP + Inches(0.38), Inches(12), Inches(0.35),
         'Different institutional architectures produce a similar effect: interpretive authority accumulates at the apex of the assemblage.',
         size=11, bold=False, color=hex_rgb('#6B7280'), italic=True)

# ─── Tier labels (left side, rotated text via thin vertical label strips) ───
# We'll use small horizontal labels that sit at the left edge
label_x = Inches(0.04)
label_w = Inches(0.32)

add_text(slide, label_x, APEX_TOP + Inches(0.2), label_w, Inches(0.5),
         'APEX', size=7, bold=True, color=hex_rgb('#D97706'), align=PP_ALIGN.CENTER)

add_text(slide, label_x, MID_TOP + Inches(0.25), label_w, Inches(0.5),
         'TRANS-\nMISSION', size=6, bold=True, color=hex_rgb('#6B7280'), align=PP_ALIGN.CENTER)

add_text(slide, label_x, BASE_TOP + Inches(0.3), label_w, Inches(0.5),
         'BASE', size=7, bold=True, color=hex_rgb('#9CA3AF'), align=PP_ALIGN.CENTER)

# ─── Horizontal convergence band at apex ───
band_y = APEX_TOP - Inches(0.06)
band_h = APEX_H + Inches(0.12)
band = add_rect(slide, L_MARGIN - Inches(0.08), band_y,
                USABLE_W + Inches(0.16), band_h,
                hex_rgb('#FFFBEB'), hex_rgb('#FDE68A'),
                radius=5, bw=Pt(1.5))

# Band label on the right
add_text(slide, SLIDE_W - Inches(3.7), band_y + Inches(0.02), Inches(3.0), Inches(0.22),
         '◀ Interpretive authority accumulates here',
         size=8, bold=True, color=hex_rgb('#B45309'), align=PP_ALIGN.RIGHT)

# ─── Build each jurisdiction column ───
for ji, jur in enumerate(JURISDICTIONS):
    col_x = L_MARGIN + ji * (COL_W + COL_GAP)
    col_cx = col_x + COL_W // 2
    jcolor = jur['color']
    jlight = jur['color_light']
    jmid = jur['color_mid']
    jborder = jur['color_border']

    # ── Header: Jurisdiction name + architecture badge ──
    # Colored dot + name
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, col_cx - Inches(0.05),
                                  HEADER_TOP + Inches(0.04), Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = hex_rgb(jcolor)
    dot.line.fill.background()

    add_text(slide, col_x, HEADER_TOP + Inches(0.02), COL_W, Inches(0.2),
             jur['name'], size=11, bold=True, color=hex_rgb('#111827'), align=PP_ALIGN.CENTER)

    # Architecture type badge
    badge_w = COL_W - Inches(0.2)
    rich_cell(slide, col_x + Inches(0.1), HEADER_TOP + Inches(0.25),
              badge_w, Inches(0.28),
              hex_rgb(jlight), hex_rgb(jborder),
              [{'t': jur['arch_type'], 's': 7, 'b': False, 'c': jcolor, 'a': PP_ALIGN.CENTER}],
              bw=Pt(0.75), mt=0.02, radius=3, anchor='middle')

    # ── FUNNEL SHAPE: Trapezoid from base to apex ──
    # We'll draw 3 shapes: apex (narrow), mid (medium), base (wide)
    # and connect them with converging lines

    # Widths: apex narrowest, base widest
    apex_w = COL_W * 0.50
    mid_w = COL_W * 0.78
    base_w = COL_W * 0.96

    apex_x = col_cx - apex_w // 2
    mid_x = col_cx - mid_w // 2
    base_x = col_cx - base_w // 2

    # ── Converging guide lines (funnel silhouette) ──
    # Left edge: base_left_x → apex_left_x
    base_left = int(base_x)
    base_right = int(base_x + base_w)
    apex_left = int(apex_x)
    apex_right = int(apex_x + apex_w)

    funnel_line_top = int(APEX_TOP + APEX_H)
    funnel_line_bot = int(BASE_TOP + BASE_H)

    # Draw subtle converging lines
    add_arrow(slide, base_left, funnel_line_bot, apex_left, funnel_line_top,
              jborder, w=1.0, dashed=True, head=False)
    add_arrow(slide, base_right, funnel_line_bot, apex_right, funnel_line_top,
              jborder, w=1.0, dashed=True, head=False)

    # ── Apex node (narrow, bold) ──
    rich_cell(slide, apex_x, APEX_TOP, apex_w, APEX_H,
              hex_rgb('#FFFFFF'), hex_rgb(jcolor),
              [{'t': jur['apex'], 's': 9.5, 'b': True, 'c': jcolor, 'a': PP_ALIGN.CENTER}],
              bw=Pt(2.5), mt=0.08, radius=5, anchor='middle')

    # ── Upward arrow: mid → apex ──
    add_arrow(slide, int(col_cx), int(MID_TOP), int(col_cx), int(APEX_TOP + APEX_H),
              jcolor, w=2.5)

    # ── Mid-layer nodes ──
    mid_node_h = Inches(0.45)
    mid_gap = Inches(0.06)
    mid_total_h = len(jur['mid']) * mid_node_h + (len(jur['mid']) - 1) * mid_gap
    mid_start_y = MID_TOP + (MID_H - mid_total_h) // 2

    for mi, mtxt in enumerate(jur['mid']):
        my = mid_start_y + mi * (mid_node_h + mid_gap)
        rich_cell(slide, mid_x, my, mid_w, mid_node_h,
                  hex_rgb(jlight), hex_rgb(jborder),
                  [{'t': mtxt, 's': 8, 'b': False, 'c': '#374151', 'a': PP_ALIGN.CENTER}],
                  bw=Pt(0.75), mt=0.03, radius=3, anchor='middle')

    # ── Upward arrow: base → mid ──
    add_arrow(slide, int(col_cx), int(BASE_TOP), int(col_cx), int(MID_TOP + MID_H),
              '#9CA3AF', w=1.8, dashed=False)

    # ── Base layer nodes ──
    base_node_w = (base_w - Inches(0.06)) // 3 - Inches(0.02)
    base_node_h = BASE_H - Inches(0.08)

    for bi, btxt in enumerate(jur['base']):
        bx = base_x + Inches(0.02) + bi * (base_node_w + Inches(0.04))
        rich_cell(slide, bx, BASE_TOP + Inches(0.04), base_node_w, base_node_h,
                  hex_rgb(jlight), hex_rgb(jborder),
                  [{'t': btxt, 's': 7, 'b': False, 'c': '#6B7280', 'a': PP_ALIGN.CENTER}],
                  bw=Pt(0.5), mt=0.06, radius=3, anchor='middle')

# ──────────────────────────────────────────────────────────────
# CONVERGENCE OBSERVATION CALLOUT
# ──────────────────────────────────────────────────────────────

call_y = Inches(6.15)
call_h = Inches(0.55)
call_bg = add_rect(slide, L_MARGIN, call_y, USABLE_W, call_h,
                    hex_rgb('#FFFBEB'), hex_rgb('#FDE68A'), radius=5, bw=Pt(1.2))

# Accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L_MARGIN, call_y, Inches(0.06), call_h)
bar.fill.solid(); bar.fill.fore_color.rgb = hex_rgb('#D97706')
bar.line.fill.background()

# Callout text with mixed formatting
tf = call_bg.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.18)
tf.margin_right = Inches(0.15)
tf.margin_top = Inches(0.06)

p = tf.paragraphs[0]

run1 = p.add_run()
run1.text = 'Convergence finding: '
run1.font.size = Pt(8.5); run1.font.bold = True
run1.font.color.rgb = hex_rgb('#B45309'); run1.font.name = 'Calibri'

run2 = p.add_run()
run2.text = ('Five jurisdictions adopt fundamentally different institutional architectures — '
             'regulatory pyramid (EU), monocentric enforcement (Colorado), multi-hub coordination (India), '
             'centralized coordinator (Brazil), and distributed constellation (U.S.) — yet all five concentrate ')
run2.font.size = Pt(8.5); run2.font.color.rgb = hex_rgb('#78350F'); run2.font.name = 'Calibri'

run3 = p.add_run()
run3.text = 'interpretive authority at apex technocratic nodes.'
run3.font.size = Pt(8.5); run3.font.bold = True
run3.font.color.rgb = hex_rgb('#B45309'); run3.font.name = 'Calibri'

run4 = p.add_run()
run4.text = ' The funnel shape is structural, not contingent.'
run4.font.size = Pt(8.5); run4.font.italic = True
run4.font.color.rgb = hex_rgb('#92400E'); run4.font.name = 'Calibri'

# ──────────────────────────────────────────────────────────────
# LEGEND ROW
# ──────────────────────────────────────────────────────────────

leg_y = Inches(6.82)

# Row labels legend
tier_labels = [
    ('Apex', '#D97706',  'Interpretive authority — standard-setting, enforcement, coordination'),
    ('Transmission', '#6B7280', 'Standards bodies, procurement, conformity assessment'),
    ('Base', '#9CA3AF',  'Statutory instruments, deployer obligations, implementation plans'),
]

lx = L_MARGIN + Inches(0.2)
for label, lcolor, desc in tier_labels:
    # Colored square
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 lx, leg_y + Inches(0.04), Inches(0.14), Inches(0.14))
    sq.fill.solid(); sq.fill.fore_color.rgb = hex_rgb(lcolor)
    sq.line.fill.background()

    add_text(slide, lx + Inches(0.2), leg_y - Inches(0.01), Inches(0.8), Inches(0.2),
             label, size=8, bold=True, color=hex_rgb(lcolor))

    add_text(slide, lx + Inches(0.85), leg_y - Inches(0.01), Inches(3.2), Inches(0.2),
             '— ' + desc, size=7.5, color=hex_rgb('#6B7280'))

    lx += Inches(4.1)

# ──────────────────────────────────────────────────────────────
# JURISDICTION COLOR KEY (bottom-right)
# ──────────────────────────────────────────────────────────────

key_y = Inches(7.1)
key_x = L_MARGIN + Inches(0.2)

for ji, jur in enumerate(JURISDICTIONS):
    dx = key_x + ji * Inches(2.2)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dx, key_y + Inches(0.03),
                                  Inches(0.09), Inches(0.09))
    dot.fill.solid(); dot.fill.fore_color.rgb = hex_rgb(jur['color'])
    dot.line.fill.background()

    add_text(slide, dx + Inches(0.14), key_y - Inches(0.02), Inches(1.8), Inches(0.2),
             jur['name'], size=8, color=hex_rgb('#374151'))

# ──────────────────────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────────────────────

output = r'c:\Users\mount\.gemini\antigravity\scratch\Authority_Concentration_Apex_Funnel.pptx'
prs.save(output)
print(f'[OK] Saved: {output}')
print('All shapes, arrows, and text are native PowerPoint objects — fully editable.')
