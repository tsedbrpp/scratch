"""
TST Propositions Evidence Matrix — Publication Quality v3
Fixed: proposition column text fully visible, no badge overlap, clean layout.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, l, t, w, h, fill, border, radius=6, bw=Pt(1)):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = bw
    sp = sh._element; pg = sp.find('.//a:prstGeom', NSMAP)
    if pg is not None:
        av = pg.find('a:avLst', NSMAP)
        if av is None: av = etree.SubElement(pg, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for c in list(av): av.remove(c)
        gd = etree.SubElement(av, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj'); gd.set('fmla', f'val {radius*1000}')
    sh.text_frame.clear()
    return sh

def add_text(slide, l, t, w, h, text, size=10, bold=False, color=None, align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.bold = bold; p.font.name = font; p.alignment = align
    if color: p.font.color.rgb = color
    return tb

def rich_card(slide, x, y, w, h, bg, bd, lines, bw=Pt(1), ml=0.15, mr=0.12, mt=0.1, mb=0.06):
    card = add_rect(slide, x, y, w, h, bg, bd, radius=5, bw=bw)
    tf = card.text_frame
    tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Inches(ml); tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mb)
    first = True
    for ln in lines:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln.get('t', '')
        r.font.size = Pt(ln.get('s', 10))
        r.font.bold = ln.get('b', False)
        r.font.name = 'Calibri'
        r.font.italic = ln.get('i', False)
        if 'c' in ln: r.font.color.rgb = hex_rgb(ln['c'])
        p.space_before = Pt(ln.get('sb', 2))
        p.space_after = Pt(ln.get('sa', 2))
        p.alignment = ln.get('a', PP_ALIGN.LEFT)
    return card

# ══════════════════════════════════════════════════════════════
# DATA
# ══════════════════════════════════════════════════════════════

JURISDICTIONS = ['EU AI Act', 'Brazil PL 2338', 'India NITI Aayog', 'Colorado SB 205']
JC = [
    ('#EFF6FF', '#93C5FD', '#1E3A8A'),
    ('#ECFDF5', '#6EE7B7', '#064E3B'),
    ('#FFF7ED', '#FDBA74', '#7C2D12'),
    ('#F5F3FF', '#C4B5FD', '#4C1D95'),
]

# Proposition data — combined lines for the prop card to control layout precisely
# badge_label combines support level into the card text itself to avoid overlapping badges
PROPS = [
    {
        'id': 'P1',
        'name': 'Referential Drift',
        'defn': 'Labels travel, referents shift\nacross jurisdictions',
        'support': 'Strong',
        'support_c': '#16A34A',
        'badge_bg': '#DCFCE7', 'badge_bd': '#86EFAC',
        'cells': [
            '"Risk" = product safety\n(Annex III enumeration)',
            '"Risk" = adaptive social\nharm (objective liability;\ncounter-drift confirmed)',
            '"Risk" = ministerial\npriorities (developmental\nresilience)',
            '"Risk" = consequential\ndecisions (consumer\nharm via AG)',
        ],
    },
    {
        'id': 'P2',
        'name': 'Infrastructural\nEmbedding',
        'defn': 'Drifted categories locked\ninto compliance systems',
        'support': 'Moderate–Strong',
        'support_c': '#D97706',
        'badge_bg': '#FEF3C7', 'badge_bd': '#FCD34D',
        'cells': [
            'CE marking, notified\nbodies, EU database\n— strong closure',
            'ANPD oversight,\nLGPD-derived\naudit routines',
            'Principles-based only;\nno binding\ninfrastructure',
            'AG rulemaking,\nduty of care;\nno private right of action',
        ],
    },
    {
        'id': 'P3',
        'name': 'Stratified Legibility',
        'defn': 'Infrastructures stabilize\ninequality of visibility',
        'support': 'Strong',
        'support_c': '#16A34A',
        'badge_bg': '#DCFCE7', 'badge_bd': '#86EFAC',
        'cells': [
            'Absorbs challenges\ntechnocratically\n(sedimentation)',
            'ANPD provides formal\ncounter-translation\n(unique)',
            'Informal pathways\nonly; no institutional\nchannel',
            'AG sole gatekeeper;\nno civil society\nstanding',
        ],
    },
]

GHOSTS = [
    'Civil Society NGOs\nWorkers / Labor\nEnd Users',
    'Indigenous Peoples\nWorkers / Labor\nCivil Society',
    'Informal Workers\nCivil Society',
    'Small Businesses\nWorkers / Unions',
]

# ══════════════════════════════════════════════════════════════
# LAYOUT — wider prop column, taller rows
# ══════════════════════════════════════════════════════════════

SLIDE_W = Inches(13.333); SLIDE_H = Inches(7.5)
prs = Presentation()
prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = hex_rgb('#FFFFFF')

GRID_L = Inches(0.3)
PROP_W = Inches(2.9)       # wider for full text visibility
GAP = Inches(0.08)
GRID_TOP = Inches(0.88)
HDR_H = Inches(0.45)
ROW_H = Inches(1.65)       # taller to fit all content
GN_H = Inches(0.85)

AVAIL = SLIDE_W - GRID_L - PROP_W - GAP - Inches(0.3)
CELL_W = (AVAIL - 3 * GAP) // 4

# ─── Title ───
add_text(slide, Inches(0.35), Inches(0.08), Inches(10), Inches(0.4),
         'TST Evidence Matrix — Three-Stage Pipeline',
         size=22, bold=True, color=hex_rgb('#111827'))

add_text(slide, Inches(0.35), Inches(0.46), Inches(12), Inches(0.3),
         'Labels travel, referents drift, infrastructures stabilize inequality',
         size=11, bold=False, color=hex_rgb('#6B7280'))

# ─── Column headers ───
add_text(slide, GRID_L + Inches(0.06), GRID_TOP + Inches(0.1), PROP_W - Inches(0.12), Inches(0.25),
         'PROPOSITION', size=9, bold=True, color=hex_rgb('#94A3B8'))

for ji, (jname, (jbg, jbd, jtxt)) in enumerate(zip(JURISDICTIONS, JC)):
    jx = GRID_L + PROP_W + GAP + ji * (CELL_W + GAP)
    rich_card(slide, jx, GRID_TOP, CELL_W, HDR_H,
              hex_rgb(jbg), hex_rgb(jbd), [
                  {'t': jname, 's': 13, 'b': True, 'c': jtxt, 'a': PP_ALIGN.CENTER}
              ], bw=Pt(1.5), mt=0.1)

# ─── Proposition rows ───
for pi, prop in enumerate(PROPS):
    ry = GRID_TOP + HDR_H + GAP + pi * (ROW_H + GAP)

    # Left: proposition card — use separate text boxes for precise control
    prop_card = add_rect(slide, GRID_L, ry, PROP_W, ROW_H,
                         hex_rgb('#F8FAFC'), hex_rgb('#CBD5E1'), radius=5, bw=Pt(1.5))

    # P-number
    add_text(slide, GRID_L + Inches(0.12), ry + Inches(0.08),
             PROP_W - Inches(0.24), Inches(0.3),
             prop['id'], size=20, bold=True, color=hex_rgb('#0F172A'))

    # Proposition name
    add_text(slide, GRID_L + Inches(0.12), ry + Inches(0.38),
             PROP_W - Inches(0.24), Inches(0.5),
             prop['name'], size=13, bold=True, color=hex_rgb('#1E293B'))

    # Definition (italic, smaller)
    add_text(slide, GRID_L + Inches(0.12), ry + Inches(0.82),
             PROP_W - Inches(0.24), Inches(0.4),
             prop['defn'], size=9, bold=False, color=hex_rgb('#64748B'))

    # Support badge — positioned at bottom of card, no overlap
    badge = add_rect(slide, GRID_L + Inches(0.12), ry + ROW_H - Inches(0.32),
                     Inches(1.6), Inches(0.24),
                     hex_rgb(prop['badge_bg']), hex_rgb(prop['badge_bd']), radius=4, bw=Pt(1))
    tf = badge.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.05); tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f'● {prop["support"]}'
    r.font.size = Pt(10); r.font.bold = True; r.font.name = 'Calibri'
    r.font.color.rgb = hex_rgb(prop['support_c'])
    p.alignment = PP_ALIGN.CENTER

    # Jurisdiction cells
    for ji, cell_text in enumerate(prop['cells']):
        jbg, jbd, jtxt = JC[ji]
        cx = GRID_L + PROP_W + GAP + ji * (CELL_W + GAP)

        cell_lines = []
        for li, line in enumerate(cell_text.split('\n')):
            cell_lines.append({
                't': line, 's': 12 if li == 0 else 11, 'b': (li == 0), 'c': jtxt,
                'sb': 0 if li == 0 else 2, 'sa': 1,
            })

        rich_card(slide, cx, ry, CELL_W, ROW_H,
                  hex_rgb('#FFFFFF'), hex_rgb(jbd), cell_lines, ml=0.12, mt=0.12)

# ─── Ghost Node Row ───
gn_y = GRID_TOP + HDR_H + GAP + 3 * (ROW_H + GAP)

# Ghost node left card
gn_card = add_rect(slide, GRID_L, gn_y, PROP_W, GN_H,
                   hex_rgb('#FEF2F2'), hex_rgb('#FECACA'), radius=5, bw=Pt(1.5))
add_text(slide, GRID_L + Inches(0.12), gn_y + Inches(0.08),
         PROP_W - Inches(0.24), Inches(0.32),
         'Ghost Nodes', size=16, bold=True, color=hex_rgb('#991B1B'))
add_text(slide, GRID_L + Inches(0.12), gn_y + Inches(0.4),
         PROP_W - Inches(0.24), Inches(0.3),
         'Excluded actors (GNDP v1.0)', size=10, bold=False, color=hex_rgb('#B91C1C'))

for ji, gn_text in enumerate(GHOSTS):
    gx = GRID_L + PROP_W + GAP + ji * (CELL_W + GAP)
    gn_lines = []
    for li, line in enumerate(gn_text.split('\n')):
        gn_lines.append({
            't': line, 's': 12 if li == 0 else 11, 'b': (li == 0),
            'c': '#991B1B' if li == 0 else '#7F1D1D',
            'sb': 0 if li == 0 else 2, 'sa': 1,
        })
    rich_card(slide, gx, gn_y, CELL_W, GN_H,
              hex_rgb('#FEF2F2'), hex_rgb('#FECACA'), gn_lines, ml=0.12, mt=0.1)

# ─── Convergence callout ───
cy = gn_y + GN_H + Inches(0.1)
cw = SLIDE_W - 2 * GRID_L
callout = add_rect(slide, GRID_L, cy, cw, Inches(0.48),
                   hex_rgb('#EEF2FF'), hex_rgb('#A5B4FC'), radius=5, bw=Pt(1.2))

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, GRID_L, cy, Inches(0.06), Inches(0.48))
bar.fill.solid(); bar.fill.fore_color.rgb = hex_rgb('#4F46E5'); bar.line.fill.background()

tf = callout.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.08)

p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = 'Convergence finding: '
r1.font.size = Pt(11); r1.font.bold = True; r1.font.name = 'Calibri'
r1.font.color.rgb = hex_rgb('#312E81')

r2 = p.add_run()
r2.text = ('All four jurisdictions deploy the same governance vocabulary but construct different governance objects. '
           'The same actor categories — civil society, workers, end users — are structurally excluded across all regimes.')
r2.font.size = Pt(10); r2.font.bold = False; r2.font.name = 'Calibri'
r2.font.color.rgb = hex_rgb('#3730A3')

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════

output = r'c:\Users\mount\.gemini\antigravity\scratch\TST_Propositions_Matrix.pptx'
prs.save(output)
print(f'[OK] Saved: {output}')
