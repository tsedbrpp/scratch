"""
Governance Styles – Technocratic Core vs. Legitimacy Framing
Matrix diagram showing the shared technocratic core across five AI governance regimes
and each regime's distinct legitimacy framing.
All elements are native PowerPoint shapes — fully editable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

NSMAP = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# ──────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────

def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_rect(slide, left, top, width, height, fill, border, radius=6, bw=Pt(1)):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border; sh.line.width = bw
    # adjust corner radius via XML
    sp = sh._element
    pg = sp.find('.//a:prstGeom', NSMAP)
    if pg is not None:
        av = pg.find('a:avLst', NSMAP)
        if av is None:
            av = etree.SubElement(pg, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for c in list(av): av.remove(c)
        gd = etree.SubElement(av, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj'); gd.set('fmla', f'val {radius * 1000}')
    sh.text_frame.clear()
    return sh

def add_text(slide, left, top, width, height, text, size=10, bold=False, color=None,
             align=PP_ALIGN.LEFT, font='Inter', italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.bold = bold; p.font.name = font; p.alignment = align; p.font.italic = italic
    if color: p.font.color.rgb = color
    return tb

def rich_cell(slide, x, y, w, h, bg, bd, lines, bw=Pt(1), ml=0.1, mt=0.08, radius=5, anchor='top'):
    card = add_rect(slide, x, y, w, h, bg, bd, radius=radius, bw=bw)
    tf = card.text_frame
    tf.word_wrap = True; tf.auto_size = None
    tf.margin_left = Inches(ml); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(0.04)
    if anchor == 'middle':
        bodyPr = tf._txBody.find('{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
    first = True
    for ln in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = ln.get('t', '')
        r.font.size = Pt(ln.get('s', 10))
        r.font.bold = ln.get('b', False)
        r.font.italic = ln.get('i', False)
        r.font.name = 'Inter'
        if 'c' in ln:
            r.font.color.rgb = hex_rgb(ln['c'])
        p.space_before = Pt(ln.get('sb', 1))
        p.space_after = Pt(ln.get('sa', 1))
        p.alignment = ln.get('a', PP_ALIGN.LEFT)
    return card

def add_arrow(slide, x1, y1, x2, y2, color_hex, w=2.0, dashed=False, head=True):
    cxn = slide.shapes._spTree.makeelement('{http://schemas.openxmlformats.org/presentationml/2006/main}cxnSp', {})
    nv = etree.SubElement(cxn, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvCxnSpPr')
    cp = etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
    cp.set('id', str(200 + len(slide.shapes)))
    cp.set('name', f'Arr{len(slide.shapes)}')
    etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvCxnSpPr')
    etree.SubElement(nv, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvPr')
    sp = etree.SubElement(cxn, '{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
    xf = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
    if x2 < x1:
        xf.set('flipH', '1')
    if y2 < y1:
        xf.set('flipV', '1')
    off = etree.SubElement(xf, '{http://schemas.openxmlformats.org/drawingml/2006/main}off')
    off.set('x', str(int(min(x1, x2))))
    off.set('y', str(int(min(y1, y2))))
    ext = etree.SubElement(xf, '{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
    ext.set('cx', str(int(abs(x2 - x1))))
    ext.set('cy', str(int(abs(y2 - y1))))
    pg = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    pg.set('prst', 'line')
    etree.SubElement(pg, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
    c = color_hex.lstrip('#')
    ln = etree.SubElement(sp, '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    ln.set('w', str(int(w * 12700)))
    sf = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    sr = etree.SubElement(sf, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
    sr.set('val', c.upper())
    if dashed:
        pd = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
        pd.set('val', 'dash')
    if head:
        te = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
        te.set('type', 'triangle'); te.set('w', 'lg'); te.set('len', 'lg')
    slide.shapes._spTree.append(cxn)

# ──────────────────────────────────────────────────────────────
# DATA
# ──────────────────────────────────────────────────────────────

JURISDICTIONS = [
    {
        'name': 'European Union',
        'color': '#6366F1',
        'light': '#EEF2FF',
        'border': '#A5B4FC',
        'legitimacy': 'Market‑harmonisation\n& supply‑chain compliance',
        'anchor': 'European Commission\n+ AI Office',
    },
    {
        'name': 'Colorado',
        'color': '#16A34A',
        'light': '#F0FDF4',
        'border': '#86EFAC',
        'legitimacy': 'Consumer‑protection\n& AG enforcement',
        'anchor': 'Attorney General',
    },
    {
        'name': 'India',
        'color': '#D97706',
        'light': '#FFFBEB',
        'border': '#FCD34D',
        'legitimacy': 'Mission‑oriented\nstate capacity',
        'anchor': 'CARO + AIREC\n+ IndiaAI Mission',
    },
    {
        'name': 'Brazil',
        'color': '#0D9488',
        'light': '#F0FDFA',
        'border': '#5EEAD4',
        'legitimacy': 'Civic‑developmental\nlanguage & ANPD oversight',
        'anchor': 'ANPD',
    },
    {
        'name': 'United States',
        'color': '#6B7280',
        'light': '#F9FAFB',
        'border': '#9CA3AF',
        'legitimacy': 'Security‑industrial\n& federal logic',
        'anchor': 'OMB + CAISI\n+ CAIOC',
    },
]

# ──────────────────────────────────────────────────────────────
# LAYOUT CONSTANTS
# ──────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

L_MARGIN = Inches(0.35)
R_MARGIN = Inches(0.35)
USABLE_W = SLIDE_W - L_MARGIN - R_MARGIN

# Column widths
JUR_W = Inches(1.4)          # jurisdiction name column
LEG_W = Inches(3.2)          # legitimacy framing column (left of centre)
CENT_W = Inches(3.0)         # central technocratic band
ANCH_W = Inches(3.2)         # institutional anchor column (right of centre)
DIV_W = Inches(0.06)         # thin divider column

# Vertical layout
TITLE_TOP = Inches(0.12)
SUB_TOP = TITLE_TOP + Inches(0.38)
HEADER_TOP = Inches(0.92)
HEADER_H = Inches(0.38)
ROW_GAP = Inches(0.06)
ROW_H = Inches(1.02)
ROW1_TOP = HEADER_TOP + HEADER_H + ROW_GAP

# X positions
JUR_X = L_MARGIN
LEG_X = JUR_X + JUR_W + Inches(0.08)
DIV_X = LEG_X + LEG_W + DIV_W
CENT_X = DIV_X + DIV_W + Inches(0.08)
ANCH_X = CENT_X + CENT_W + DIV_W + Inches(0.08)

# ──────────────────────────────────────────────────────────────
# BUILD
# ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = hex_rgb('#FFFFFF')

# Title
add_text(slide, Inches(0.4), TITLE_TOP, Inches(12), Inches(0.4),
         'Governance Styles – Technocratic Core vs. Legitimacy Framing',
         size=20, bold=True, color=hex_rgb('#111827'))
add_text(slide, Inches(0.4), SUB_TOP, Inches(12), Inches(0.35),
         'All regimes share a compliance‑driven technocratic core, but each legitimises it differently.',
         size=10.5, bold=False, color=hex_rgb('#6B7280'), italic=True)

# Header row (column titles)
# Jurisdiction header (just a thin background for alignment)
rich_cell(slide, JUR_X, HEADER_TOP, JUR_W, HEADER_H,
          hex_rgb('#1E293B'), hex_rgb('#0F172A'),
          [{'t': 'JURISDICTION', 's': 8, 'b': True, 'c': '#F8FAFC', 'a': PP_ALIGN.CENTER}],
          bw=Pt(0), mt=0.06, anchor='middle')
# Legitimacy framing header
rich_cell(slide, LEG_X, HEADER_TOP, LEG_W, HEADER_H,
          hex_rgb('#FFFBEB'), hex_rgb('#FDE68A'),
          [{'t': 'LEGITIMACY FRAMING', 's': 8, 'b': True, 'c': '#92400E', 'a': PP_ALIGN.CENTER}],
          bw=Pt(1), mt=0.04, anchor='middle')
# Central technocratic core header
rich_cell(slide, CENT_X, HEADER_TOP, CENT_W, HEADER_H,
          hex_rgb('#F9FAFB'), hex_rgb('#D1D5DB'),
          [{'t': 'TECHNOCRATIC CORE', 's': 8, 'b': True, 'c': '#374151', 'a': PP_ALIGN.CENTER}],
          bw=Pt(1), mt=0.04, anchor='middle')
# Institutional anchor header
rich_cell(slide, ANCH_X, HEADER_TOP, ANCH_W, HEADER_H,
          hex_rgb('#FEF2F2'), hex_rgb('#FECACA'),
          [{'t': 'INSTITUTIONAL ANCHOR', 's': 8, 'b': True, 'c': '#991B1B', 'a': PP_ALIGN.CENTER}],
          bw=Pt(1), mt=0.04, anchor='middle')

# Central technocratic band (visual representation of core)
core_bg = add_rect(slide, CENT_X, ROW1_TOP - Inches(0.02), CENT_W, 5 * ROW_H + 4 * ROW_GAP + Inches(0.04),
                   hex_rgb('#F9FAFB'), hex_rgb('#D1D5DB'), radius=6, bw=Pt(1.5))
# Inside the band, three stacked icons/text (using simple text blocks)
icon_y = ROW1_TOP + ROW_H * 0.2
add_text(slide, CENT_X + Inches(0.12), icon_y, CENT_W - Inches(0.24), Inches(0.3),
         'Standards', size=10, bold=True, color=hex_rgb('#374151'), align=PP_ALIGN.CENTER)
add_text(slide, CENT_X + Inches(0.12), icon_y + Inches(0.4), CENT_W - Inches(0.24), Inches(0.3),
         'Compliance', size=10, bold=True, color=hex_rgb('#374151'), align=PP_ALIGN.CENTER)
add_text(slide, CENT_X + Inches(0.12), icon_y + Inches(0.8), CENT_W - Inches(0.24), Inches(0.3),
         'Administration', size=10, bold=True, color=hex_rgb('#374151'), align=PP_ALIGN.CENTER)

# Rows for each jurisdiction
for i, jur in enumerate(JURISDICTIONS):
    row_y = ROW1_TOP + i * (ROW_H + ROW_GAP)
    # Jurisdiction name cell
    rich_cell(slide, JUR_X, row_y, JUR_W, ROW_H,
              hex_rgb(jur['light']), hex_rgb(jur['border']),
              [{'t': jur['name'], 's': 11, 'b': True, 'c': jur['color'], 'a': PP_ALIGN.CENTER}],
              bw=Pt(1.5), mt=0.06, anchor='middle')
    # Legitimacy framing cell
    rich_cell(slide, LEG_X, row_y, LEG_W, ROW_H,
              hex_rgb(jur['light']), hex_rgb(jur['border']),
              [{'t': jur['legitimacy'], 's': 9, 'b': False, 'c': '#78350F', 'a': PP_ALIGN.CENTER}],
              bw=Pt(1), mt=0.08, anchor='middle')
    # Institutional anchor cell
    rich_cell(slide, ANCH_X, row_y, ANCH_W, ROW_H,
              hex_rgb(jur['light']), hex_rgb(jur['border']),
              [{'t': jur['anchor'], 's': 9, 'b': False, 'c': jur['color'], 'a': PP_ALIGN.CENTER}],
              bw=Pt(1), mt=0.08, anchor='middle')
    # Dashed arrows from core band to both outer cells
    # From centre of core band to centre of legitimacy cell
    add_arrow(slide, int(CENT_X + CENT_W/2), int(row_y + ROW_H/2),
              int(LEG_X + LEG_W/2), int(row_y + ROW_H/2),
              jur['color'], w=1.5, dashed=True, head=False)
    # From centre of core band to centre of anchor cell
    add_arrow(slide, int(CENT_X + CENT_W/2), int(row_y + ROW_H/2),
              int(ANCH_X + ANCH_W/2), int(row_y + ROW_H/2),
              jur['color'], w=1.5, dashed=True, head=False)

# Bottom callout
call_y = ROW1_TOP + 5 * (ROW_H + ROW_GAP) + Inches(0.08)
call_h = Inches(0.55)
call_bg = add_rect(slide, L_MARGIN, call_y, USABLE_W, call_h,
                    hex_rgb('#FEF2F2'), hex_rgb('#FECACA'), radius=5, bw=Pt(1.2))
# Accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L_MARGIN, call_y, Inches(0.06), call_h)
bar.fill.solid(); bar.fill.fore_color.rgb = hex_rgb('#991B1B')
bar.line.fill.background()
# Text inside callout
tf = call_bg.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.06)
para = tf.paragraphs[0]
run1 = para.add_run()
run1.text = 'Shared core: '
run1.font.size = Pt(8.5); run1.font.bold = True; run1.font.color.rgb = hex_rgb('#991B1B')
run2 = para.add_run()
run2.text = 'compliance‑driven technocratic architecture'
run2.font.size = Pt(8.5); run2.font.color.rgb = hex_rgb('#7F1D1D')
run3 = para.add_run()
run3.text = '. Distinct legitimacy framings (market, consumer, mission, civic, security) sit around this core.'
run3.font.size = Pt(8.5); run3.font.italic = True; run3.font.color.rgb = hex_rgb('#991B1B')

# Legend row
leg_y = call_y + call_h + Inches(0.14)
# Core legend
add_text(slide, L_MARGIN + Inches(0.3), leg_y, Inches(2.5), Inches(0.2),
         'Technocratic core (standards → compliance → administration)', size=8, color=hex_rgb('#374151'))
# Legitimacy legend
add_text(slide, L_MARGIN + Inches(4.0), leg_y, Inches(2.5), Inches(0.2),
         'Legitimacy framing (market, consumer, mission, civic, security)', size=8, color=hex_rgb('#78350F'))
# Anchor legend
add_text(slide, L_MARGIN + Inches(7.5), leg_y, Inches(2.5), Inches(0.2),
         'Institutional anchor (commission, AG, mission bodies, ANPD, OMB/CAISI)', size=8, color=hex_rgb('#991B1B'))

# Save
output = r'c:\Users\mount\.gemini\antigravity\scratch\Governance_Styles_Matrix.pptx'
prs.save(output)
print(f'[OK] Saved: {output}')
print('All shapes, arrows, and text are native PowerPoint objects — fully editable.')
